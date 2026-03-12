"""pptx-ea-font 기반 한글 폰트 후처리.

최종 PPTX 파일의 모든 텍스트에 EA(East Asian) 폰트 속성을 일괄 적용합니다.
pandoc 생성 슬라이드와 python-pptx 커스텀 슬라이드 모두에 균일하게 적용됩니다.
"""

from __future__ import annotations

import logging
from pathlib import Path

from pptx import Presentation
from pptx_ea_font import set_font as set_ea_font

from ppt_maker.errors import FontError
from ppt_maker.fonts.resolver import ResolvedFont, resolve_font

logger = logging.getLogger(__name__)


def apply_korean_font(
    pptx_path: Path,
    font: ResolvedFont | None = None,
    output_path: Path | None = None,
) -> Path:
    """PPTX 파일의 모든 텍스트에 한글(EA) 폰트를 적용.

    Args:
        pptx_path: 원본 PPTX 파일 경로.
        font: 적용할 폰트. None이면 자동 탐색.
        output_path: 저장 경로. None이면 원본 덮어쓰기.

    Returns:
        저장된 PPTX 파일 경로.
    """
    if not pptx_path.exists():
        raise FontError(f"PPTX 파일을 찾을 수 없습니다: {pptx_path}")

    if font is None:
        font = resolve_font()

    output_path = output_path or pptx_path

    try:
        prs = Presentation(str(pptx_path))
        count = 0

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            set_ea_font(run, font.name)
                            count += 1

                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            for paragraph in cell.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    set_ea_font(run, font.name)
                                    count += 1

        prs.save(str(output_path))
        logger.info("한글 폰트 '%s' 적용 완료: %d개 텍스트 요소", font.name, count)
        return output_path

    except Exception as e:
        raise FontError(
            f"한글 폰트 적용 중 오류 발생: {e}",
            detail=f"파일: {pptx_path}, 폰트: {font.name}",
        ) from e
