"""python-pptx 커스텀 슬라이드 빌더.

구조화된 데이터에서 커스텀 슬라이드를 생성하는 빌더 패턴 구현.
템플릿 플레이스홀더를 우선 사용하고, 없으면 직접 도형 배치.
"""

from __future__ import annotations

import logging
from pathlib import Path
from typing import TYPE_CHECKING

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from ppt_maker.errors import PptxCustomError
from ppt_maker.theme.palette import ThemeConfig

if TYPE_CHECKING:
    from ppt_maker.assets.manager import AssetManager

logger = logging.getLogger(__name__)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


# --- 헬퍼 함수 ---


def set_bg(slide, color: RGBColor) -> None:
    """슬라이드 배경색 설정."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(
    slide,
    left: int | float,
    top: int | float,
    width: int | float,
    height: int | float,
    fill_color: RGBColor | None = None,
    border_color: RGBColor | None = None,
    border_width: int = 0,
) -> object:
    """사각형 도형 추가."""
    from pptx.enum.shapes import MSO_SHAPE

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()

    if border_color and border_width:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width)
    else:
        shape.line.fill.background()

    return shape


def add_text(
    slide,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    font_size: int = 18,
    font_color: RGBColor = RGBColor(0xFF, 0xFF, 0xFF),
    font_name: str = "Apple SD Gothic Neo",
    bold: bool = False,
    alignment: PP_ALIGN = PP_ALIGN.LEFT,
) -> object:
    """텍스트 상자 추가."""
    tx_box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    tf = tx_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.name = font_name
    p.font.bold = bold
    p.alignment = alignment
    return tx_box


def add_bullet_list(
    slide,
    items: list[str],
    left: float,
    top: float,
    width: float,
    height: float,
    *,
    font_size: int = 16,
    font_color: RGBColor = RGBColor(0xFF, 0xFF, 0xFF),
    font_name: str = "Apple SD Gothic Neo",
    bullet_color: RGBColor | None = None,
) -> object:
    """불릿 리스트 텍스트 상자 추가."""
    tx_box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    tf = tx_box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.name = font_name
        p.space_after = Pt(6)

    return tx_box


def _parse_md_runs(text: str) -> list[tuple[str, bool]]:
    """마크다운 **bold** → (텍스트, is_bold) 튜플 리스트."""
    import re
    parts: list[tuple[str, bool]] = []
    last_end = 0
    for m in re.finditer(r"\*\*(.+?)\*\*", text):
        if m.start() > last_end:
            parts.append((text[last_end:m.start()], False))
        parts.append((m.group(1), True))
        last_end = m.end()
    if last_end < len(text):
        parts.append((text[last_end:], False))
    return parts or [("", False)]


def _set_paragraph_with_bold(p, text: str) -> None:
    """paragraph에 마크다운 bold를 실제 pptx bold run으로 변환."""
    runs = _parse_md_runs(text)
    p.clear()
    for run_text, is_bold in runs:
        run = p.add_run()
        run.text = run_text
        run.font.bold = is_bold


def _fill_placeholder(slide, ph_idx: int, text: str) -> bool:
    """플레이스홀더에 텍스트를 채운다. **bold** → 실제 bold."""
    try:
        ph = slide.placeholders[ph_idx]
        tf = ph.text_frame
        tf.clear()
        _set_paragraph_with_bold(tf.paragraphs[0], text)
        return True
    except (KeyError, IndexError):
        return False


def _fill_placeholder_bullets(slide, ph_idx: int, items: list[str]) -> bool:
    """플레이스홀더에 불릿 리스트를 채운다. **bold** → 실제 bold."""
    try:
        ph = slide.placeholders[ph_idx]
        tf = ph.text_frame
        tf.clear()
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            _set_paragraph_with_bold(p, item)
            p.space_after = Pt(4)
        return True
    except (KeyError, IndexError):
        return False


def _fill_placeholder_content(slide, ph_idx: int, data: dict) -> bool:
    """items가 있으면 불릿으로, 없으면 content 텍스트로 플레이스홀더를 채운다."""
    items = data.get("items", [])
    if items:
        return _fill_placeholder_bullets(slide, ph_idx, items)
    content = data.get("content", "")
    if content:
        return _fill_placeholder(slide, ph_idx, content)
    return False


# --- 슬라이드 타입별 빌더 ---


class SlideBuilder:
    """커스텀 슬라이드 빌더.

    회사 템플릿이 등록되어 있으면 플레이스홀더를 우선 채우고,
    없으면 Blank 레이아웃에 도형을 직접 배치합니다.
    """

    def __init__(
        self,
        theme: ThemeConfig,
        layout_mapping: dict[str, int] | None = None,
        asset_manager: AssetManager | None = None,
    ) -> None:
        self.theme = theme
        self.colors = theme.colors
        self.font_name = theme.fonts.heading
        self._layout_mapping = layout_mapping or {}
        self.asset_manager = asset_manager

    @property
    def _has_template(self) -> bool:
        """템플릿 레이아웃 매핑이 있는지 확인."""
        return bool(self._layout_mapping)

    # 매핑 없는 타입의 폴백 체인
    _LAYOUT_FALLBACKS: dict[str, list[str]] = {
        "keynote": ["content", "section", "blank"],
        "three_content": ["two_content", "content", "blank"],
        "grid_2x2": ["two_content", "content", "blank"],
        "grid_3": ["two_content", "content", "blank"],
        "content_area": ["content", "blank"],
        "content_area_dark": ["content", "blank"],
        "picture_left": ["two_content", "content", "blank"],
        "picture_right": ["two_content", "content", "blank"],
    }

    def _get_layout(self, prs: Presentation, slide_type: str) -> object:
        """슬라이드 타입에 매핑된 레이아웃 반환. 폴백 체인 탐색."""
        idx = self._layout_mapping.get(slide_type)
        if idx is not None and idx < len(prs.slide_layouts):
            return prs.slide_layouts[idx]
        # 폴백 체인에서 매핑 찾기
        for fallback_type in self._LAYOUT_FALLBACKS.get(slide_type, []):
            fb_idx = self._layout_mapping.get(fallback_type)
            if fb_idx is not None and fb_idx < len(prs.slide_layouts):
                return prs.slide_layouts[fb_idx]
        # 최종 폴백: Blank
        blank_idx = self._layout_mapping.get("blank")
        if blank_idx is not None and blank_idx < len(prs.slide_layouts):
            return prs.slide_layouts[blank_idx]
        fallback = min(6, len(prs.slide_layouts) - 1)
        return prs.slide_layouts[fallback]

    def _add_slide(self, prs: Presentation, slide_type: str) -> object:
        """레이아웃 매핑 기반 슬라이드 추가."""
        layout = self._get_layout(prs, slide_type)
        return prs.slides.add_slide(layout)

    # --- 빌더 메서드 ---

    def build_title(self, prs: Presentation, data: dict) -> None:
        """타이틀 슬라이드. 템플릿: ph0=CENTER_TITLE, ph1=SUBTITLE."""
        slide = self._add_slide(prs, "title")

        title = data.get("title", "")
        subtitle = data.get("subtitle", data.get("content", ""))

        if self._has_template:
            _fill_placeholder(slide, 0, title)
            if subtitle:
                _fill_placeholder(slide, 1, subtitle)
        else:
            set_bg(slide, self.colors.bg_primary)
            add_text(
                slide, title,
                1.0, 2.0, 11.0, 1.5,
                font_size=40, font_color=self.colors.text_primary,
                font_name=self.font_name, bold=True,
                alignment=PP_ALIGN.CENTER,
            )
            if subtitle:
                add_text(
                    slide, subtitle,
                    1.0, 3.8, 11.0, 1.0,
                    font_size=24, font_color=self.colors.accent,
                    font_name=self.font_name,
                    alignment=PP_ALIGN.CENTER,
                )

    def build_content(self, prs: Presentation, data: dict) -> None:
        """일반 콘텐츠 슬라이드. 템플릿: ph0=TITLE, ph1=OBJECT(본문)."""
        slide = self._add_slide(prs, "content")

        title = data.get("title", "")
        items = data.get("items", [])
        content = data.get("content", "")

        if self._has_template:
            _fill_placeholder(slide, 0, title)
            if items:
                _fill_placeholder_bullets(slide, 1, items)
            elif content:
                _fill_placeholder(slide, 1, content)
        else:
            set_bg(slide, self.colors.bg_primary)
            add_text(
                slide, title,
                0.5, 0.3, 12.0, 0.8,
                font_size=28, font_color=self.colors.text_primary,
                font_name=self.font_name, bold=True,
            )
            if items:
                add_bullet_list(
                    slide, items,
                    0.8, 1.5, 11.5, 5.5,
                    font_size=16, font_color=self.colors.text_primary,
                    font_name=self.font_name,
                )
            elif content:
                add_text(
                    slide, content,
                    0.8, 1.5, 11.5, 5.5,
                    font_size=16, font_color=self.colors.text_primary,
                    font_name=self.font_name,
                )

    def build_section(self, prs: Presentation, data: dict) -> None:
        """섹션 구분 슬라이드. 템플릿: ph0=TITLE, ph1=BODY."""
        slide = self._add_slide(prs, "section")

        title = data.get("title", "")
        subtitle = data.get("subtitle", data.get("content", ""))

        if self._has_template:
            _fill_placeholder(slide, 0, title)
            if subtitle:
                _fill_placeholder(slide, 1, subtitle)
        else:
            set_bg(slide, self.colors.bg_primary)
            add_text(
                slide, title,
                1.0, 2.5, 11.0, 1.5,
                font_size=36, font_color=self.colors.text_primary,
                font_name=self.font_name, bold=True,
                alignment=PP_ALIGN.LEFT,
            )
            if subtitle:
                add_text(
                    slide, subtitle,
                    1.0, 4.2, 11.0, 1.0,
                    font_size=20, font_color=self.colors.text_secondary,
                    font_name=self.font_name,
                    alignment=PP_ALIGN.LEFT,
                )

    def build_comparison(self, prs: Presentation, data: dict) -> None:
        """좌우 비교 슬라이드. 구조화 데이터 없으면 content 폴백."""
        left = data.get("left", {})
        right = data.get("right", {})

        # 구조화 데이터가 없으면 items를 반반 나눠서 비교 슬라이드로 표시
        if not left and not right:
            items = data.get("items", [])
            content = data.get("content", "")
            if items and len(items) >= 2:
                mid = len(items) // 2
                left = {"title": "", "items": items[:mid]}
                right = {"title": "", "items": items[mid:]}
            elif content:
                # content로 폴백
                self.build_content(prs, data)
                return

        slide = self._add_slide(prs, "comparison")
        title = data.get("title", "")

        if self._has_template:
            _fill_placeholder(slide, 0, title)
            # ph1: left label, ph2: left content, ph3: right label, ph4: right content
            _fill_placeholder(slide, 1, left.get("title", ""))
            left_items = left.get("items", [])
            if left_items:
                _fill_placeholder_bullets(slide, 2, left_items)
            elif left.get("content"):
                _fill_placeholder(slide, 2, left["content"])
            _fill_placeholder(slide, 3, right.get("title", ""))
            right_items = right.get("items", [])
            if right_items:
                _fill_placeholder_bullets(slide, 4, right_items)
            elif right.get("content"):
                _fill_placeholder(slide, 4, right["content"])
        else:
            set_bg(slide, self.colors.bg_primary)
            add_text(
                slide, title,
                0.5, 0.3, 12.0, 0.8,
                font_size=28, font_color=self.colors.text_primary,
                font_name=self.font_name, bold=True,
            )
            add_shape(slide, 0.5, 1.5, 5.8, 5.5, fill_color=self.colors.bg_secondary)
            add_text(
                slide, left.get("title", ""),
                0.8, 1.7, 5.2, 0.6,
                font_size=22, font_color=self.colors.accent,
                font_name=self.font_name, bold=True,
            )
            if left.get("items"):
                add_bullet_list(
                    slide, left["items"],
                    0.8, 2.5, 5.2, 4.0,
                    font_color=self.colors.text_primary,
                    font_name=self.font_name,
                )
            add_shape(slide, 6.8, 1.5, 5.8, 5.5, fill_color=self.colors.bg_secondary)
            add_text(
                slide, right.get("title", ""),
                7.1, 1.7, 5.2, 0.6,
                font_size=22, font_color=self.colors.accent2,
                font_name=self.font_name, bold=True,
            )
            if right.get("items"):
                add_bullet_list(
                    slide, right["items"],
                    7.1, 2.5, 5.2, 4.0,
                    font_color=self.colors.text_primary,
                    font_name=self.font_name,
                )

    def build_timeline(self, prs: Presentation, data: dict) -> None:
        """타임라인 슬라이드. 구조화 events가 없으면 content 폴백."""
        events = data.get("events", [])

        if not events:
            # content 폴백: Title and Content 레이아웃 사용
            self.build_content(prs, data)
            return

        slide = self._add_slide(prs, "three_content")
        title = data.get("title", "Timeline")

        if self._has_template:
            _fill_placeholder(slide, 0, title)
        else:
            set_bg(slide, self.colors.bg_primary)
            add_text(
                slide, title,
                0.5, 0.3, 12.0, 0.8,
                font_size=28, font_color=self.colors.text_primary,
                font_name=self.font_name, bold=True,
            )

        add_shape(slide, 1.0, 3.5, 11.0, 0.03, fill_color=self.colors.accent)
        accent_colors = [self.colors.accent, self.colors.accent2, self.colors.accent3]
        spacing = min(11.0 / max(len(events), 1), 3.0)

        for i, event in enumerate(events):
            x = 1.0 + i * spacing
            color = event.get("color") or accent_colors[i % len(accent_colors)]
            if isinstance(color, str):
                from ppt_maker.theme.palette import hex_to_rgb
                color = hex_to_rgb(color)

            add_shape(slide, x + spacing / 2 - 0.1, 3.35, 0.2, 0.2, fill_color=color)
            add_text(
                slide, event.get("date", ""),
                x, 2.3, spacing, 0.5,
                font_size=12, font_color=self.colors.accent,
                font_name=self.font_name, alignment=PP_ALIGN.CENTER,
            )
            add_text(
                slide, event.get("title", ""),
                x, 4.0, spacing, 1.0,
                font_size=14, font_color=self.colors.text_primary,
                font_name=self.font_name, alignment=PP_ALIGN.CENTER,
            )

    def build_keynote(self, prs: Presentation, data: dict) -> None:
        """키노트 슬라이드. message 키 없으면 content/items로 폴백."""
        slide = self._add_slide(prs, "keynote")
        title = data.get("title", "")
        message = data.get("message", "")

        # message가 없으면 content에서 추출
        if not message:
            items = data.get("items", [])
            content = data.get("content", "")
            if items:
                message = "\n".join(f"• {item}" for item in items)
            elif content:
                message = content

        if self._has_template:
            _fill_placeholder(slide, 0, title)
            # items가 있으면 불릿으로, 없으면 단일 텍스트로
            items = data.get("items", [])
            if items:
                _fill_placeholder_bullets(slide, 1, items)
            elif message:
                _fill_placeholder(slide, 1, message)
        else:
            set_bg(slide, self.colors.bg_primary)
            add_text(
                slide, title,
                1.0, 2.0, 11.0, 2.0,
                font_size=36, font_color=self.colors.text_primary,
                font_name=self.font_name, bold=True,
                alignment=PP_ALIGN.CENTER,
            )
            if message:
                add_text(
                    slide, message,
                    1.5, 4.5, 10.0, 1.5,
                    font_size=20, font_color=self.colors.text_secondary,
                    font_name=self.font_name,
                    alignment=PP_ALIGN.CENTER,
                )

    def build_quote(self, prs: Presentation, data: dict) -> None:
        """인용/질문 슬라이드."""
        slide = self._add_slide(prs, "keynote")
        quote = data.get("quote", data.get("content", ""))

        if self._has_template:
            _fill_placeholder(slide, 0, quote)
            if data.get("attribution"):
                _fill_placeholder(slide, 1, f"— {data['attribution']}")
        else:
            set_bg(slide, self.colors.bg_primary)
            add_text(
                slide, '"',
                1.0, 1.0, 1.0, 1.5,
                font_size=72, font_color=self.colors.accent,
                font_name=self.font_name, bold=True,
            )
            add_text(
                slide, quote,
                1.5, 2.5, 10.0, 2.0,
                font_size=28, font_color=self.colors.text_primary,
                font_name=self.font_name,
                alignment=PP_ALIGN.CENTER,
            )
            if data.get("attribution"):
                add_text(
                    slide, f"— {data['attribution']}",
                    1.5, 5.0, 10.0, 0.5,
                    font_size=16, font_color=self.colors.text_secondary,
                    font_name=self.font_name,
                    alignment=PP_ALIGN.RIGHT,
                )

    def build_card_list(self, prs: Presentation, data: dict) -> None:
        """카드형 목록 슬라이드."""
        slide = self._add_slide(prs, "grid_2x2")

        title = data.get("title", "")
        if self._has_template:
            _fill_placeholder(slide, 0, title)
        else:
            set_bg(slide, self.colors.bg_primary)
            add_text(
                slide, title,
                0.5, 0.3, 12.0, 0.8,
                font_size=28, font_color=self.colors.text_primary,
                font_name=self.font_name, bold=True,
            )

        items = data.get("items", [])
        accent_colors = [self.colors.accent, self.colors.accent2, self.colors.accent3]

        for i, item in enumerate(items[:6]):
            col = i % 3
            row = i // 3
            x = 0.5 + col * 4.2
            y = 1.5 + row * 3.0

            add_shape(slide, x, y, 3.8, 2.5, fill_color=self.colors.bg_secondary)
            item_title = item.get("title", item) if isinstance(item, dict) else str(item)
            add_text(
                slide, item_title,
                x + 0.2, y + 0.2, 3.4, 0.5,
                font_size=18, font_color=accent_colors[i % len(accent_colors)],
                font_name=self.font_name, bold=True,
            )
            if isinstance(item, dict) and item.get("detail"):
                add_text(
                    slide, item["detail"],
                    x + 0.2, y + 0.8, 3.4, 1.5,
                    font_size=14, font_color=self.colors.text_primary,
                    font_name=self.font_name,
                )

    def build_process(self, prs: Presentation, data: dict) -> None:
        """단계별 프로세스 슬라이드. steps가 없으면 content 폴백."""
        steps = data.get("steps", [])
        if not steps:
            self.build_content(prs, data)
            return

        slide = self._add_slide(prs, "grid_3")
        title = data.get("title", "Process")

        if self._has_template:
            _fill_placeholder(slide, 0, title)
        else:
            set_bg(slide, self.colors.bg_primary)
            add_text(
                slide, title,
                0.5, 0.3, 12.0, 0.8,
                font_size=28, font_color=self.colors.text_primary,
                font_name=self.font_name, bold=True,
            )

        accent_colors = [self.colors.accent, self.colors.accent2, self.colors.accent3]
        n = len(steps)
        step_width = min(12.0 / n, 3.0)
        start_x = (13.333 - step_width * n) / 2

        for i, step in enumerate(steps):
            x = start_x + i * step_width
            color = accent_colors[i % len(accent_colors)]

            add_shape(slide, x + step_width / 2 - 0.3, 2.0, 0.6, 0.6, fill_color=color)
            add_text(
                slide, str(i + 1),
                x + step_width / 2 - 0.3, 2.05, 0.6, 0.6,
                font_size=20, font_color=self.colors.bg_primary,
                font_name=self.font_name, bold=True,
                alignment=PP_ALIGN.CENTER,
            )
            add_text(
                slide, step.get("label", ""),
                x, 2.8, step_width, 0.5,
                font_size=16, font_color=self.colors.text_primary,
                font_name=self.font_name, bold=True,
                alignment=PP_ALIGN.CENTER,
            )
            if step.get("desc"):
                add_text(
                    slide, step["desc"],
                    x, 3.4, step_width, 2.0,
                    font_size=12, font_color=self.colors.text_secondary,
                    font_name=self.font_name,
                    alignment=PP_ALIGN.CENTER,
                )
            if i < n - 1:
                add_text(
                    slide, "→",
                    x + step_width - 0.2, 2.1, 0.4, 0.5,
                    font_size=24, font_color=self.colors.text_secondary,
                    font_name=self.font_name,
                    alignment=PP_ALIGN.CENTER,
                )

    def build_picture_left(self, prs: Presentation, data: dict) -> None:
        """좌측 이미지 + 우측 텍스트 슬라이드."""
        slide = self._add_slide(prs, "picture_left")

        if self._has_template:
            _fill_placeholder(slide, 0, data.get("title", ""))
            _fill_placeholder_content(slide, 1, data)
        else:
            set_bg(slide, self.colors.bg_primary)
            add_text(
                slide, data.get("title", ""),
                0.5, 0.3, 12.0, 0.8,
                font_size=28, font_color=self.colors.text_primary,
                font_name=self.font_name, bold=True,
            )

        brand_name = data.get("brand", data.get("image", ""))
        image_path: Path | None = None
        if self.asset_manager and brand_name:
            image_path = self.asset_manager.find_brand_image(brand_name)

        if image_path and image_path.exists():
            slide.shapes.add_picture(
                str(image_path),
                Inches(0.5), Inches(1.5), Inches(4.5), Inches(5.0),
            )
        elif not self._has_template:
            add_shape(
                slide, 0.5, 1.5, 5.5, 5.5,
                fill_color=self.colors.bg_secondary,
            )
            add_text(
                slide, data.get("image_label", "[Image]"),
                1.5, 3.5, 3.5, 1.0,
                font_size=16, font_color=self.colors.text_secondary,
                font_name=self.font_name,
                alignment=PP_ALIGN.CENTER,
            )

        if not self._has_template:
            if data.get("content"):
                add_text(
                    slide, data["content"],
                    6.5, 1.8, 6.0, 5.0,
                    font_size=16, font_color=self.colors.text_primary,
                    font_name=self.font_name,
                )
            if data.get("items"):
                add_bullet_list(
                    slide, data["items"],
                    6.5, 1.8, 6.0, 5.0,
                    font_color=self.colors.text_primary,
                    font_name=self.font_name,
                )

    def build_picture_right(self, prs: Presentation, data: dict) -> None:
        """좌측 텍스트 + 우측 이미지 슬라이드."""
        slide = self._add_slide(prs, "picture_right")

        if self._has_template:
            _fill_placeholder(slide, 0, data.get("title", ""))
            _fill_placeholder_content(slide, 1, data)
        else:
            set_bg(slide, self.colors.bg_primary)
            add_text(
                slide, data.get("title", ""),
                0.5, 0.3, 12.0, 0.8,
                font_size=28, font_color=self.colors.text_primary,
                font_name=self.font_name, bold=True,
            )

        brand_name = data.get("brand", data.get("image", ""))
        image_path: Path | None = None
        if self.asset_manager and brand_name:
            image_path = self.asset_manager.find_brand_image(brand_name)

        if image_path and image_path.exists():
            slide.shapes.add_picture(
                str(image_path),
                Inches(8.0), Inches(1.5), Inches(4.5), Inches(5.0),
            )
        elif not self._has_template:
            add_shape(
                slide, 7.0, 1.5, 5.5, 5.5,
                fill_color=self.colors.bg_secondary,
            )
            add_text(
                slide, data.get("image_label", "[Image]"),
                8.0, 3.5, 3.5, 1.0,
                font_size=16, font_color=self.colors.text_secondary,
                font_name=self.font_name,
                alignment=PP_ALIGN.CENTER,
            )

        if not self._has_template:
            if data.get("content"):
                add_text(
                    slide, data["content"],
                    0.5, 1.8, 6.0, 5.0,
                    font_size=16, font_color=self.colors.text_primary,
                    font_name=self.font_name,
                )
            if data.get("items"):
                add_bullet_list(
                    slide, data["items"],
                    0.5, 1.8, 6.0, 5.0,
                    font_color=self.colors.text_primary,
                    font_name=self.font_name,
                )

    def build_blank(self, prs: Presentation, data: dict) -> None:
        """빈 콘텐츠 영역 슬라이드 (밝은 배경)."""
        slide = self._add_slide(prs, "content_area")
        title = data.get("title", "")
        if title and not (
            self._has_template and _fill_placeholder(slide, 0, title)
        ):
            add_text(
                slide, title,
                0.5, 0.3, 12.0, 0.8,
                font_size=28, font_color=self.colors.text_primary,
                font_name=self.font_name, bold=True,
            )

    def build_blank_dark(self, prs: Presentation, data: dict) -> None:
        """빈 콘텐츠 영역 슬라이드 (어두운 배경)."""
        slide = self._add_slide(prs, "content_area_dark")
        if not self._has_template:
            set_bg(slide, self.colors.bg_primary)
        title = data.get("title", "")
        if title and not (
            self._has_template and _fill_placeholder(slide, 0, title)
        ):
            add_text(
                slide, title,
                0.5, 0.3, 12.0, 0.8,
                font_size=28, font_color=self.colors.text_primary,
                font_name=self.font_name, bold=True,
            )

    # --- 디스패처 ---

    BUILDERS: dict[str, str] = {
        "title": "build_title",
        "section": "build_section",
        "content": "build_content",
        "quote": "build_quote",
        "timeline": "build_timeline",
        "comparison": "build_comparison",
        "card_list": "build_card_list",
        "process": "build_process",
        "picture_left": "build_picture_left",
        "picture_right": "build_picture_right",
        "blank": "build_blank",
        "blank_dark": "build_blank_dark",
        "keynote": "build_keynote",
    }

    def build_slide(self, prs: Presentation, slide_type: str, data: dict) -> None:
        """슬라이드 타입에 따라 적절한 빌더 호출."""
        method_name = self.BUILDERS.get(slide_type)
        if method_name is None:
            raise PptxCustomError(
                f"지원되지 않는 슬라이드 타입: {slide_type}",
                detail=f"지원 타입: {', '.join(sorted(self.BUILDERS))}",
            )
        method = getattr(self, method_name)
        try:
            method(prs, data)
        except PptxCustomError:
            raise
        except Exception as e:
            raise PptxCustomError(
                f"슬라이드 생성 실패 ({slide_type}): {e}",
            ) from e

    def create_presentation(
        self,
        template_path: Path | None = None,
    ) -> Presentation:
        """프레젠테이션 생성. 회사 템플릿이 있으면 해당 파일 기반."""
        if template_path and template_path.exists():
            prs = Presentation(str(template_path))
            xml_slides = prs.slides._sldIdLst
            for slide_id in list(xml_slides):
                rel_id = slide_id.get(
                    "{http://schemas.openxmlformats.org/officeDocument"
                    "/2006/relationships}id"
                )
                if rel_id:
                    prs.part.drop_rel(rel_id)
                xml_slides.remove(slide_id)
            logger.info("회사 템플릿 기반 프레젠테이션 생성: %s", template_path)
            return prs
        prs = Presentation()
        prs.slide_width = Emu(int(self.theme.slide.width * 914400))
        prs.slide_height = Emu(int(self.theme.slide.height * 914400))
        return prs
