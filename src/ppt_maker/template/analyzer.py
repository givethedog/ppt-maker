"""회사 PPTX 템플릿 분석기.

임의의 .pptx 파일에서 슬라이드 레이아웃, 플레이스홀더, 색상, 폰트를 자동 추출합니다.
"""

from __future__ import annotations

import logging
import sys
from dataclasses import dataclass, field
from pathlib import Path

if sys.version_info >= (3, 11):
    import tomllib
else:
    import tomli as tomllib

from pptx import Presentation

from ppt_maker.errors import ConfigError

logger = logging.getLogger(__name__)


@dataclass
class PlaceholderInfo:
    """플레이스홀더 정보."""

    idx: int
    name: str
    type: str  # "TITLE", "BODY", "SUBTITLE", "PICTURE", etc.
    left: float  # inches
    top: float
    width: float
    height: float


@dataclass
class LayoutInfo:
    """슬라이드 레이아웃 정보."""

    name: str
    index: int
    placeholders: list[PlaceholderInfo] = field(default_factory=list)


@dataclass
class TemplateManifest:
    """템플릿 분석 결과 매니페스트."""

    source_path: str
    slide_width: float  # inches
    slide_height: float
    layouts: list[LayoutInfo] = field(default_factory=list)
    fonts_used: set[str] = field(default_factory=set)
    layout_mapping: dict[str, int] = field(default_factory=dict)

    def get_layout_index(self, ppt_maker_type: str) -> int | None:
        """ppt-maker 슬라이드 타입에 매핑된 레이아웃 인덱스 반환."""
        return self.layout_mapping.get(ppt_maker_type)

    def summary(self) -> str:
        """분석 결과 요약 문자열."""
        lines = [
            f"템플릿: {self.source_path}",
            f"슬라이드 크기: {self.slide_width:.1f} x {self.slide_height:.1f} inches",
            f"레이아웃: {len(self.layouts)}개",
            f"감지된 폰트: {', '.join(sorted(self.fonts_used)) or '없음'}",
            "",
            "레이아웃 목록:",
        ]
        for layout in self.layouts:
            ph_str = ", ".join(f"{p.name}({p.type})" for p in layout.placeholders)
            lines.append(f"  [{layout.index}] {layout.name}: {ph_str or '플레이스홀더 없음'}")

        if self.layout_mapping:
            lines.append("")
            lines.append("슬라이드 타입 매핑:")
            for k, v in sorted(self.layout_mapping.items()):
                lines.append(f"  {k} → [{v}] {self.layouts[v].name}")

        return "\n".join(lines)


def _emu_to_inches(emu: int) -> float:
    return round(emu / 914400, 3)


def _detect_placeholder_type(ph) -> str:
    """python-pptx placeholder 타입을 문자열로 변환."""
    from pptx.enum.shapes import PP_PLACEHOLDER

    type_map = {
        PP_PLACEHOLDER.TITLE: "TITLE",
        PP_PLACEHOLDER.CENTER_TITLE: "TITLE",
        PP_PLACEHOLDER.SUBTITLE: "SUBTITLE",
        PP_PLACEHOLDER.BODY: "BODY",
        PP_PLACEHOLDER.OBJECT: "BODY",
    }
    try:
        return type_map.get(ph.placeholder_format.type, "OTHER")
    except Exception:
        return "OTHER"


# 레이아웃 이름 → ppt-maker 슬라이드 타입 자동 매핑 규칙
# BMW CI Template 등 회사 템플릿의 레이아웃 이름도 매칭
_LAYOUT_NAME_HINTS: dict[str, list[str]] = {
    "title": [
        "title slide", "제목 슬라이드", "title | full area",
        "title", "표지",
    ],
    "section": [
        "section header", "섹션 머리글", "divider", "구역",
    ],
    "content": [
        "title and content", "제목 및 내용",
        "content | 1", "본문",
    ],
    "two_content": [
        "two content", "2개의 내용", "content | 2",
    ],
    "three_content": ["content | 3"],
    "four_content": ["content | 4"],
    "comparison": ["content | 2x2", "comparison", "비교"],
    "grid_1": ["grid | 1"],
    "grid_2": ["grid | 2"],
    "grid_3": ["grid | 3"],
    "grid_4": ["grid | 4"],
    "grid_2x2": ["grid | 2x2"],
    "content_area": ["content | area"],
    "content_area_dark": ["content | area dark"],
    "picture_left": ["content | picture left", "picture left"],
    "picture_right": ["content | picture right", "picture right"],
    "picture_2": ["content | picture | 2"],
    "picture_3": ["content | picture | 3"],
    "picture_4": ["content | picture | 4"],
    "keynote": ["key note", "keynote"],
    "blank": ["blank", "빈 화면", "빈 슬라이드"],
}


def _auto_map_layout(layout_name: str) -> str | None:
    """레이아웃 이름에서 ppt-maker 슬라이드 타입 추론."""
    name_lower = layout_name.lower().strip()
    # 긴 힌트부터 매칭하여 "title"이 "title and content"보다 먼저 매칭되는 것을 방지
    candidates: list[tuple[int, str, str]] = []
    for ppt_type, hints in _LAYOUT_NAME_HINTS.items():
        for hint in hints:
            if hint in name_lower:
                candidates.append((len(hint), hint, ppt_type))
    if candidates:
        candidates.sort(reverse=True)
        return candidates[0][2]
    return None


def analyze_template(pptx_path: Path) -> TemplateManifest:
    """PPTX 템플릿을 분석하여 매니페스트 생성.

    Args:
        pptx_path: 분석할 .pptx 파일 경로.

    Returns:
        TemplateManifest: 분석 결과.
    """
    pptx_path = Path(pptx_path)
    if not pptx_path.exists():
        raise ConfigError(f"템플릿 파일을 찾을 수 없습니다: {pptx_path}")

    try:
        prs = Presentation(str(pptx_path))
    except Exception as e:
        raise ConfigError(
            f"PPTX 템플릿을 열 수 없습니다: {pptx_path}",
            detail=str(e),
        ) from e

    manifest = TemplateManifest(
        source_path=str(pptx_path),
        slide_width=_emu_to_inches(prs.slide_width),
        slide_height=_emu_to_inches(prs.slide_height),
    )

    # 레이아웃 분석
    for idx, layout in enumerate(prs.slide_layouts):
        layout_info = LayoutInfo(name=layout.name, index=idx)

        for ph in layout.placeholders:
            ph_info = PlaceholderInfo(
                idx=ph.placeholder_format.idx,
                name=ph.name,
                type=_detect_placeholder_type(ph),
                left=_emu_to_inches(ph.left),
                top=_emu_to_inches(ph.top),
                width=_emu_to_inches(ph.width),
                height=_emu_to_inches(ph.height),
            )
            layout_info.placeholders.append(ph_info)

        manifest.layouts.append(layout_info)

        # 자동 매핑
        mapped_type = _auto_map_layout(layout.name)
        if mapped_type and mapped_type not in manifest.layout_mapping:
            manifest.layout_mapping[mapped_type] = idx

    # 폰트 수집 (기존 슬라이드에서)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.name:
                            manifest.fonts_used.add(run.font.name)

    n_layouts = len(manifest.layouts)
    n_fonts = len(manifest.fonts_used)
    logger.info("템플릿 분석 완료: %d 레이아웃, %d 폰트", n_layouts, n_fonts)
    return manifest


def save_manifest(manifest: TemplateManifest, output_path: Path) -> Path:
    """매니페스트를 TOML 파일로 저장."""
    import tomli_w

    data = {
        "template": {
            "source_path": manifest.source_path,
            "slide_width": manifest.slide_width,
            "slide_height": manifest.slide_height,
            "fonts_used": sorted(manifest.fonts_used),
        },
        "layout_mapping": manifest.layout_mapping,
        "layouts": [
            {
                "name": lay.name,
                "index": lay.index,
                "placeholders": [
                    {
                        "idx": p.idx,
                        "name": p.name,
                        "type": p.type,
                        "left": p.left,
                        "top": p.top,
                        "width": p.width,
                        "height": p.height,
                    }
                    for p in lay.placeholders
                ],
            }
            for lay in manifest.layouts
        ],
    }

    output_path.parent.mkdir(parents=True, exist_ok=True)
    with open(output_path, "wb") as f:
        tomli_w.dump(data, f)

    logger.info("매니페스트 저장: %s", output_path)
    return output_path


def load_manifest(manifest_path: Path) -> TemplateManifest:
    """TOML 매니페스트 파일에서 TemplateManifest를 로드."""
    manifest_path = Path(manifest_path)
    if not manifest_path.exists():
        raise ConfigError(f"매니페스트 파일을 찾을 수 없습니다: {manifest_path}")

    with open(manifest_path, "rb") as f:
        data = tomllib.load(f)

    tpl = data.get("template", {})
    manifest = TemplateManifest(
        source_path=tpl.get("source_path", ""),
        slide_width=tpl.get("slide_width", 0),
        slide_height=tpl.get("slide_height", 0),
        fonts_used=set(tpl.get("fonts_used", [])),
        layout_mapping=data.get("layout_mapping", {}),
    )

    for layout_data in data.get("layouts", []):
        layout_info = LayoutInfo(
            name=layout_data["name"],
            index=layout_data["index"],
        )
        for ph_data in layout_data.get("placeholders", []):
            layout_info.placeholders.append(
                PlaceholderInfo(
                    idx=ph_data["idx"],
                    name=ph_data["name"],
                    type=ph_data["type"],
                    left=ph_data["left"],
                    top=ph_data["top"],
                    width=ph_data["width"],
                    height=ph_data["height"],
                )
            )
        manifest.layouts.append(layout_info)

    return manifest
