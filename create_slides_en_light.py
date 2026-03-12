#!/usr/bin/env python3
"""AI Trends 2025-2026 Presentation Slides (English Version - Light Theme)"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Light Theme Color Palette ──
BG_LIGHT = RGBColor(0xFA, 0xFA, 0xFC)      # near-white background
BG_CARD = RGBColor(0xF0, 0xF2, 0xF5)       # card background (light gray)
ACCENT = RGBColor(0x00, 0x66, 0xFF)        # deep blue (primary accent)
ACCENT2 = RGBColor(0x6C, 0x5C, 0xE7)      # indigo purple
ACCENT3 = RGBColor(0x00, 0xB8, 0x94)      # teal green
TEXT_DARK = RGBColor(0x1A, 0x1A, 0x2E)     # dark text (near black)
TEXT_MID = RGBColor(0x4A, 0x5A, 0x6A)      # mid text
TEXT_LIGHT = RGBColor(0x8A, 0x8A, 0x9A)    # light text
ORANGE = RGBColor(0xF5, 0x5A, 0x2C)       # orange
YELLOW = RGBColor(0xF0, 0xA5, 0x00)       # gold
N8N_COLOR = RGBColor(0xEA, 0x4B, 0x71)    # n8n brand (kept)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# ── Brand Colors ──
OPENAI_BLACK = RGBColor(0x00, 0x00, 0x00)
ANTHROPIC_BROWN = RGBColor(0xD4, 0xA5, 0x74)
GOOGLE_BLUE = RGBColor(0x42, 0x85, 0xF4)
OPENSOURCE_ORANGE = RGBColor(0xFF, 0x66, 0x00)
CHALLENGER_PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
LANGGRAPH_COLOR = RGBColor(0x6C, 0x5C, 0xE7)
CREWAI_COLOR = RGBColor(0x00, 0xB8, 0x94)
OPENAI_SDK_COLOR = RGBColor(0x00, 0x66, 0xFF)
CLAUDE_SDK_COLOR = RGBColor(0xF0, 0xA5, 0x00)

FONT_NAME = "Calibri"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_shape_with_border(slide, left, top, width, height, color, border_color=None):
    """Shape with a subtle border for light backgrounds"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text(slide, left, top, width, height, text, size=18, color=TEXT_DARK, bold=False, align=PP_ALIGN.LEFT, font_name=FONT_NAME):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_bullet_list(slide, left, top, width, height, items, size=16, color=TEXT_DARK, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = FONT_NAME
        p.space_after = spacing
    return txBox


def add_brand_circle(slide, left, top, size, bg_color, text, text_color=WHITE, font_size=12):
    """Brand-colored circle with initial text for visual logo representation"""
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    circle.fill.solid()
    circle.fill.fore_color.rgb = bg_color
    circle.line.fill.background()
    circle.shadow.inherit = False
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = text_color
    p.font.bold = True
    p.font.name = FONT_NAME
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return circle


# ════════════════════════════════════════════
# Slide 1: Title
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_bg(slide, BG_LIGHT)

# Top accent bar
bar_top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.12))
bar_top.fill.solid()
bar_top.fill.fore_color.rgb = ACCENT
bar_top.line.fill.background()

add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
         "AI Trends 2025-2026", size=48, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(2.8), Inches(11), Inches(1),
         "The Age of Agentic AI and the Promise of n8n", size=28, color=ACCENT, align=PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.5),
         "IT Team Knowledge Sharing  |  March 2026", size=18, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

# Bottom decorative line
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3), Inches(4.0), Inches(7), Pt(3))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

# ════════════════════════════════════════════
# Slide 2: Opening Question
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_LIGHT)

add_text(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.5),
         '"What has happened in the AI world\nover the past year?"', size=40, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, Inches(2), Inches(4.5), Inches(9), Inches(1),
         "Two years after ChatGPT... the world has changed faster than anyone imagined",
         size=20, color=TEXT_MID, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# Slide 3: Key Terminology
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_LIGHT)

add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
         "Key Terms You'll Hear Today", size=32, color=TEXT_DARK, bold=True)

add_text(slide, Inches(0.5), Inches(1.0), Inches(12), Inches(0.4),
         "Lots of AI/IT jargon ahead. Let's cover the essentials first!",
         size=16, color=TEXT_MID)

terms = [
    ("LLM", "Large Language Model", "Large-scale AI models like ChatGPT that understand and generate text", ACCENT),
    ("Agent", "AI Agent", "An AI program that reasons and acts autonomously toward goals", ACCENT2),
    ("MCP", "Model Context Protocol", "A universal standard for AI to connect to external tools (the USB-C of AI)", ACCENT3),
    ("CLI", "Command Line Interface", "Terminal/command-line environment for text-based computer interaction", YELLOW),
    ("Token", "Token", "The smallest unit AI uses to process text (roughly 1 word per token)", ORANGE),
    ("RAG", "Retrieval Augmented Generation", "AI retrieves external documents to enrich its responses", N8N_COLOR),
    ("Workflow", "Workflow", "An automated task sequence: A then B then C, executed in order", ACCENT),
    ("Orchestration", "Orchestration", "Coordinating multiple AI systems and tools, like a conductor", ACCENT2),
    ("Fair-code", "Fair-code", "Similar to open source but with restrictions on commercial redistribution", ACCENT3),
    ("No-code", "No-code/Low-code", "Building software via drag-and-drop with little or no code", N8N_COLOR),
]

for i, (term, term_full, desc, color) in enumerate(terms):
    row = i // 2
    col = i % 2
    x = Inches(0.5 + col * 6.3)
    y = Inches(1.6 + row * 1.1)

    add_shape_with_border(slide, x, y, Inches(6.0), Inches(0.95), BG_CARD, TEXT_LIGHT)

    # Left color bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.1), Inches(0.95))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    add_text(slide, x + Inches(0.3), y + Inches(0.05), Inches(2.5), Inches(0.45),
             f"{term} ({term_full})", size=15, color=color, bold=True)
    add_text(slide, x + Inches(0.3), y + Inches(0.5), Inches(5.4), Inches(0.6),
             desc, size=13, color=TEXT_MID)

add_text(slide, Inches(0.5), Inches(7.0), Inches(12), Inches(0.4),
         "Feel free to ask about any unfamiliar terms at any point during the presentation!",
         size=14, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# Slide 4: Timeline
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_LIGHT)

add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
         "The Past Year: A Timeline", size=32, color=TEXT_DARK, bold=True)

timeline = [
    ("2024.11", "MCP Announced", ACCENT),
    ("2025.01", "DeepSeek R1 Shock", ORANGE),
    ("2025.02", "Vibe Coding Coined", YELLOW),
    ("2025.03", "OpenAI Agents SDK", ACCENT2),
    ("2025.05", "Claude Code Launched", ACCENT),
    ("2025.09", "Notion AI Agents", ACCENT3),
    ("2025.10", "n8n $180M Funding", N8N_COLOR),
    ("2025.12", "MCP to Linux Fdn", TEXT_LIGHT),
    ("2026.01", "OpenClaw 250K Stars", ORANGE),
    ("2026.02", "Claude 4.6 Released", ACCENT),
]

for i, (date, event, color) in enumerate(timeline):
    row = i // 5
    col = i % 5
    x = Inches(0.4 + col * 2.5)
    y = Inches(1.5 + row * 2.8)

    # Circle point
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.9), y, Inches(0.25), Inches(0.25))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()

    add_text(slide, x, y + Inches(0.4), Inches(2.2), Inches(0.4),
             date, size=13, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x, y + Inches(0.75), Inches(2.2), Inches(0.5),
             event, size=14, color=TEXT_DARK, align=PP_ALIGN.CENTER)

# Bottom message
add_text(slide, Inches(0.5), Inches(6.7), Inches(12), Inches(0.5),
         'AI is evolving: "Conversation Tool" -> "Working Agent" -> "A Team That Runs Systems"',
         size=16, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# Slide 5: LLM Model Competition
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_LIGHT)

add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
         "The LLM Race: A Five-Way Battle", size=32, color=TEXT_DARK, bold=True)

players = [
    ("OpenAI", "GPT-4.1, o3, o4-mini", "Reasoning-focused, 1M tokens", ACCENT3, OPENAI_BLACK, "AI", WHITE),
    ("Anthropic", "Claude 4.6", "#1 in coding, 30h+ autonomous agent", ACCENT, ANTHROPIC_BROWN, "A", WHITE),
    ("Google", "Gemini 2.5 Pro", "1M tokens, Deep Think mode", YELLOW, GOOGLE_BLUE, "G", WHITE),
    ("Open Source", "DeepSeek, Llama 4", "Price disruption, 10M tokens", ORANGE, OPENSOURCE_ORANGE, "OS", WHITE),
    ("Challengers", "Grok 3, Mistral", "200K GPU infra + EU data sovereignty", ACCENT2, CHALLENGER_PURPLE, "X", WHITE),
]

for i, (name, models, strength, color, brand_bg, brand_text, brand_text_color) in enumerate(players):
    y = Inches(1.3 + i * 1.15)
    card = add_shape_with_border(slide, Inches(0.5), y, Inches(12), Inches(1.0), BG_CARD, TEXT_LIGHT)

    # Left color bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), y, Inches(0.12), Inches(1.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    # Brand color circle icon
    add_brand_circle(slide, Inches(0.8), y + Inches(0.15), Inches(0.7), brand_bg, brand_text, brand_text_color, font_size=14)

    add_text(slide, Inches(1.7), y + Inches(0.1), Inches(2), Inches(0.5),
             name, size=20, color=color, bold=True)
    add_text(slide, Inches(3.8), y + Inches(0.1), Inches(3.5), Inches(0.5),
             models, size=16, color=TEXT_DARK)
    add_text(slide, Inches(7.5), y + Inches(0.1), Inches(5), Inches(0.5),
             strength, size=16, color=TEXT_MID)

# Bottom
add_text(slide, Inches(0.5), Inches(6.8), Inches(12), Inches(0.5),
         "The competitive axis: Benchmarks -> Agent capability -> Team orchestration",
         size=15, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# Slide 6: Competitive Axis Shift
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_LIGHT)

add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
         "How the Competitive Axis Has Shifted", size=32, color=TEXT_DARK, bold=True)

phases = [
    ("2024", "Competing on benchmark scores", TEXT_LIGHT),
    ("2025", "How long can it work autonomously?", ACCENT),
    ("2026", "How well can it orchestrate agent teams?", ACCENT2),
]
for i, (year, desc, color) in enumerate(phases):
    y = Inches(1.5 + i * 1.6)
    add_text(slide, Inches(1.5), y, Inches(2), Inches(0.6),
             year, size=36, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(4), y + Inches(0.05), Inches(7), Inches(0.6),
             desc, size=22, color=TEXT_DARK)
    if i < 2:
        add_text(slide, Inches(2.2), y + Inches(0.8), Inches(1), Inches(0.5),
                 "v", size=30, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

trends = [
    "Every model now ships with reasoning (thinking) mode -- table stakes",
    "Across-the-board price drops since DeepSeek R1 -- cost barrier disappearing",
    "Open-source models rapidly closing the gap with proprietary ones",
]
add_bullet_list(slide, Inches(1), Inches(5.8), Inches(11), Inches(1.5),
                trends, size=15, color=TEXT_MID)

# ════════════════════════════════════════════
# Slide 7: MCP Overview
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_LIGHT)

add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
         "MCP: The USB-C of AI", size=32, color=TEXT_DARK, bold=True)

add_text(slide, Inches(0.5), Inches(1.2), Inches(12), Inches(0.6),
         "Model Context Protocol -- a universal standard for AI to connect with external tools",
         size=18, color=TEXT_MID)

# Before
add_shape_with_border(slide, Inches(0.5), Inches(2.2), Inches(5.8), Inches(3.5), BG_CARD, TEXT_LIGHT)
add_text(slide, Inches(0.8), Inches(2.3), Inches(5), Inches(0.5),
         "Before MCP", size=20, color=ORANGE, bold=True)
before_items = [
    "AI Model A --- custom code --- GitHub",
    "AI Model A --- custom code --- Slack",
    "AI Model B --- custom code --- GitHub",
    "...",
    "M x N integrations needed",
]
add_bullet_list(slide, Inches(0.8), Inches(2.9), Inches(5), Inches(2.5),
                before_items, size=14, color=TEXT_MID)

# After
add_shape_with_border(slide, Inches(7), Inches(2.2), Inches(5.8), Inches(3.5), BG_CARD, TEXT_LIGHT)
add_text(slide, Inches(7.3), Inches(2.3), Inches(5), Inches(0.5),
         "After MCP", size=20, color=ACCENT3, bold=True)
after_items = [
    "AI Model A --+",
    "AI Model B --+-- MCP --+-- GitHub Server",
    "AI Model C --+         +-- Slack Server",
    "                       +-- DB Server",
    "Only M + N implementations needed!",
]
add_bullet_list(slide, Inches(7.3), Inches(2.9), Inches(5), Inches(2.5),
                after_items, size=14, color=TEXT_MID)

# Stats
stats = [
    ("5,800+", "MCP Servers"),
    ("300+", "Clients"),
    ("97M+", "Monthly Downloads"),
]
for i, (num, label) in enumerate(stats):
    x = Inches(1.5 + i * 4)
    add_text(slide, x, Inches(6.0), Inches(3), Inches(0.5),
             num, size=28, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x, Inches(6.5), Inches(3), Inches(0.4),
             label, size=14, color=TEXT_MID, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# Slide 8: MCP Limitations
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_LIGHT)

add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
         "But Wait... Is MCP Already Showing Its Age?", size=32, color=ORANGE, bold=True)

add_text(slide, Inches(0.5), Inches(1.2), Inches(12), Inches(0.6),
         "Major services are embedding their own AI agents natively",
         size=18, color=TEXT_MID)

services = [
    ("Notion", "3.0 AI Agents -> 3.3 Custom Agents", "Native is primary, MCP alongside"),
    ("Linear", "GitHub Copilot Agent native integration", "Direct integration without MCP"),
    ("GitHub", "Copilot CLI + Agent Mode", "Native CLI more efficient than MCP"),
]

for i, (name, feature, note) in enumerate(services):
    y = Inches(2.2 + i * 1.3)
    add_shape_with_border(slide, Inches(0.5), y, Inches(12), Inches(1.1), BG_CARD, TEXT_LIGHT)
    add_text(slide, Inches(0.9), y + Inches(0.15), Inches(2), Inches(0.4),
             name, size=20, color=ACCENT, bold=True)
    add_text(slide, Inches(3.2), y + Inches(0.15), Inches(5), Inches(0.4),
             feature, size=16, color=TEXT_DARK)
    add_text(slide, Inches(8.5), y + Inches(0.15), Inches(4), Inches(0.4),
             note, size=14, color=TEXT_MID)

add_text(slide, Inches(0.5), Inches(6.0), Inches(12), Inches(0.4),
         "Major services are embedding native AI, reducing their reliance on MCP",
         size=15, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# Slide 9: CLI vs MCP Efficiency
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_LIGHT)

add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
         "CLI vs MCP: The Efficiency Debate", size=32, color=TEXT_DARK, bold=True)

add_text(slide, Inches(0.5), Inches(1.2), Inches(12), Inches(0.5),
         "CLI tools may be more efficient for AI agents than MCP",
         size=18, color=TEXT_MID)

# Comparison card - CLI
add_shape_with_border(slide, Inches(0.5), Inches(2.2), Inches(5.8), Inches(3.0), BG_CARD, TEXT_LIGHT)
add_text(slide, Inches(0.8), Inches(2.3), Inches(5.2), Inches(0.5),
         "CLI Approach (gh, kubectl, etc.)", size=20, color=ACCENT3, bold=True)
cli_items = [
    "Token cost: ~200 tokens",
    "LLMs already know CLI from training data",
    "33% efficiency advantage (benchmark)",
    "e.g., gh pr list -> executes immediately",
]
add_bullet_list(slide, Inches(0.8), Inches(3.0), Inches(5.2), Inches(2.0),
                cli_items, size=15, color=TEXT_DARK)

# Comparison card - MCP
add_shape_with_border(slide, Inches(7), Inches(2.2), Inches(5.8), Inches(3.0), BG_CARD, TEXT_LIGHT)
add_text(slide, Inches(7.3), Inches(2.3), Inches(5.2), Inches(0.5),
         "MCP Approach", size=20, color=ORANGE, bold=True)
mcp_items = [
    "Token cost: ~55,000 tokens (GitHub, 93 tools)",
    "Schemas encountered at runtime for the first time",
    "Excessive context consumption",
    "Highly versatile, but inefficient",
]
add_bullet_list(slide, Inches(7.3), Inches(3.0), Inches(5.2), Inches(2.0),
                mcp_items, size=15, color=TEXT_MID)

# Conclusion
add_shape_with_border(slide, Inches(1.5), Inches(5.8), Inches(10), Inches(1.2), BG_CARD, ACCENT)
add_text(slide, Inches(1.8), Inches(5.9), Inches(9.4), Inches(1.0),
         "MCP is not going away (it is already a standard)\n"
         "But it is shifting from \"master key\" to \"one of several approaches\"\n"
         "The real winner: the orchestration layer that combines all of the above",
         size=16, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# Slide 10: Return to CLI & Vibe Coding
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_LIGHT)

add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
         "The Return to CLI & Vibe Coding", size=32, color=TEXT_DARK, bold=True)

# Evolution stages
stages = [
    ("2023", "Autocomplete", "Copilot, Tabnine", TEXT_LIGHT),
    ("2024", "Chat-based IDE", "Cursor, Windsurf", ACCENT2),
    ("2025", "CLI Agents", "Claude Code, Aider", ACCENT),
    ("2026", "Team Agents", "Claude Teams, n8n", N8N_COLOR),
]
for i, (year, stage, tools, color) in enumerate(stages):
    x = Inches(0.5 + i * 3.2)
    add_shape_with_border(slide, x, Inches(1.5), Inches(2.8), Inches(1.8), BG_CARD, TEXT_LIGHT)
    add_text(slide, x, Inches(1.6), Inches(2.8), Inches(0.4),
             year, size=24, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x, Inches(2.1), Inches(2.8), Inches(0.4),
             stage, size=18, color=TEXT_DARK, align=PP_ALIGN.CENTER)
    add_text(slide, x, Inches(2.6), Inches(2.8), Inches(0.4),
             tools, size=13, color=TEXT_MID, align=PP_ALIGN.CENTER)
    if i < 3:
        add_text(slide, x + Inches(2.8), Inches(2.0), Inches(0.4), Inches(0.5),
                 "->", size=24, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

# Claude Code
add_text(slide, Inches(0.5), Inches(3.8), Inches(6), Inches(0.5),
         "Claude Code: #1 developer preference (46%) within 8 months", size=18, color=ACCENT, bold=True)

# Vibe Coding
add_shape_with_border(slide, Inches(0.5), Inches(4.5), Inches(6), Inches(2.5), BG_CARD, TEXT_LIGHT)
add_text(slide, Inches(0.8), Inches(4.6), Inches(5.5), Inches(0.5),
         'Vibe Coding (Karpathy, Feb 2025)', size=18, color=YELLOW, bold=True)
add_text(slide, Inches(0.8), Inches(5.1), Inches(5.5), Inches(0.8),
         '"Forget that the code even exists"\n-> OpenClaw: 250K stars, shipped without reading code\n-> 2.74x more security vulnerabilities',
         size=14, color=TEXT_MID)

add_shape_with_border(slide, Inches(7), Inches(4.5), Inches(5.8), Inches(2.5), BG_CARD, TEXT_LIGHT)
add_text(slide, Inches(7.3), Inches(4.6), Inches(5.2), Inches(0.5),
         'Agentic Engineering (2026)', size=18, color=ACCENT3, bold=True)
add_text(slide, Inches(7.3), Inches(5.1), Inches(5.2), Inches(0.8),
         '"Vibe coding is already behind us" -- Karpathy\n-> AI writes code, humans supervise and orchestrate\n-> Quality and security matter',
         size=14, color=TEXT_MID)

add_text(slide, Inches(6.2), Inches(5.3), Inches(0.6), Inches(0.5),
         "->", size=28, color=ACCENT, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# Slide 11: Context Engineering
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_LIGHT)

add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
         "From Prompt to Context Engineering", size=32, color=TEXT_DARK, bold=True)

add_text(slide, Inches(0.5), Inches(1.1), Inches(12), Inches(0.5),
         '"The new essential skill in AI is not prompting -- it\'s context engineering" -- Tobi Lutke (Shopify CEO)',
         size=16, color=TEXT_MID)

# Comparison table
headers = ["", "Prompt Engineering", "Context Engineering"]
rows = [
    ["Scope", "A single text string", "The entire information architecture"],
    ["Timing", "One-shot at invocation", "Continuous, across multiple turns"],
    ["Focus", "What should I ask?", "What does the AI already know?"],
    ["Outcome", "Better phrasing", "Better system design"],
]

for i, header in enumerate(headers):
    x = Inches(0.5 + i * 4.2)
    add_shape(slide, x, Inches(1.8), Inches(4.0), Inches(0.6), ACCENT2 if i > 0 else BG_CARD)
    add_text(slide, x + Inches(0.2), Inches(1.85), Inches(3.6), Inches(0.5),
             header, size=16, color=WHITE if i > 0 else TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)

for ri, row in enumerate(rows):
    for ci, cell in enumerate(row):
        x = Inches(0.5 + ci * 4.2)
        y = Inches(2.5 + ri * 0.65)
        bg = BG_CARD if ri % 2 == 0 else BG_LIGHT
        add_shape_with_border(slide, x, y, Inches(4.0), Inches(0.6), bg, TEXT_LIGHT if ri % 2 == 0 else None)
        color = ACCENT if ci == 2 else (TEXT_MID if ci == 1 else TEXT_DARK)
        add_text(slide, x + Inches(0.2), y + Inches(0.05), Inches(3.6), Inches(0.5),
                 cell, size=14, color=color, align=PP_ALIGN.CENTER)

add_text(slide, Inches(0.5), Inches(5.2), Inches(12), Inches(0.4),
         "Prompt = \"What should I ask?\"  vs  Context = \"What does the AI already know?\"",
         size=16, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# Slide 12: Context Engineering - 7 Elements
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_LIGHT)

add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
         "The 7 Elements of Context", size=32, color=TEXT_DARK, bold=True)

add_text(slide, Inches(0.5), Inches(1.1), Inches(12), Inches(0.5),
         "Design the information you provide to AI across seven layers",
         size=18, color=TEXT_MID)

elements_detail = [
    ("System Prompt", "AI's role and rules", '"You are a customer service expert"', ACCENT),
    ("User Prompt", "The current task/question", '"Draft a reply to this email"', ACCENT2),
    ("Conversation History", "Context so far", "References previous exchanges", ACCENT3),
    ("Long-term Memory", "Persistence across sessions", "Remembers last week's work", N8N_COLOR),
    ("Retrieved Info (RAG)", "Relevant documents/data", "Internal manuals retrieved and referenced", ORANGE),
    ("Available Tools", "APIs, MCP, functions", "Slack messaging, DB queries available", ACCENT),
    ("Output Format", "Structured response def.", "JSON, tables, Markdown, etc.", YELLOW),
]

for i, (name, desc, example, color) in enumerate(elements_detail):
    row = i // 2
    col = i % 2
    x = Inches(0.5 + col * 6.3)
    y = Inches(1.8 + row * 1.3)

    add_shape_with_border(slide, x, y, Inches(6.0), Inches(1.15), BG_CARD, TEXT_LIGHT)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.1), Inches(1.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    add_text(slide, x + Inches(0.3), y + Inches(0.05), Inches(2.5), Inches(0.4),
             name, size=16, color=color, bold=True)
    add_text(slide, x + Inches(3.0), y + Inches(0.05), Inches(2.8), Inches(0.4),
             desc, size=13, color=TEXT_DARK)
    add_text(slide, x + Inches(0.3), y + Inches(0.55), Inches(5.5), Inches(0.4),
             f"e.g., {example}", size=12, color=TEXT_MID)

add_text(slide, Inches(0.5), Inches(7.0), Inches(12), Inches(0.4),
         "Context engineering = designing what, when, in what format, and how much to provide to AI",
         size=14, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# Bridge Slide: So Far -> Agentic AI
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_LIGHT)

add_text(slide, Inches(1), Inches(1.0), Inches(11), Inches(0.8),
         "Let's recap where we are...", size=36, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)

bridge_items = [
    ("Models are smart enough", TEXT_LIGHT),
    ("Tool connectivity standards (MCP) exist", ACCENT),
    ("AI can write code on its own", ACCENT2),
    ("Context management has evolved", ACCENT3),
]

for i, (text, color) in enumerate(bridge_items):
    y = Inches(2.3 + i * 0.9)
    add_text(slide, Inches(2), y, Inches(9), Inches(0.6),
             f"   {text}", size=22, color=color)

add_text(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.8),
         "So what happens when all of this comes together?", size=28, color=TEXT_MID, align=PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(6.3), Inches(11), Inches(0.8),
         "->  Agentic AI", size=40, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# Agentic AI - Definition
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_LIGHT)

add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
         "Agentic AI: What Is It, Really?", size=32, color=TEXT_DARK, bold=True)

add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(0.8),
         "Generative AI \"answers\" questions.\nAgentic AI \"achieves\" goals.",
         size=28, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

# Restaurant analogy
add_shape_with_border(slide, Inches(0.5), Inches(3.0), Inches(5.8), Inches(3.8), BG_CARD, TEXT_LIGHT)
add_text(slide, Inches(0.8), Inches(3.1), Inches(5.2), Inches(0.5),
         "Generative AI = Chef who gives you the recipe", size=18, color=TEXT_MID, bold=True)
items_gen = [
    'Input: "Give me a pasta recipe"',
    'Output: Text (the recipe)',
    'Tools: None',
    'Next step: Human decides',
]
add_bullet_list(slide, Inches(0.8), Inches(3.7), Inches(5.2), Inches(2.5),
                items_gen, size=15, color=TEXT_MID)

add_shape_with_border(slide, Inches(7), Inches(3.0), Inches(5.8), Inches(3.8), BG_CARD, ACCENT)
add_text(slide, Inches(7.3), Inches(3.1), Inches(5.2), Inches(0.5),
         "Agentic AI = Chef who cooks and serves", size=18, color=ACCENT, bold=True)
items_agent = [
    'Input: "Prepare a dinner party"',
    'Output: Actions (shop -> cook -> set table -> serve)',
    'Tools: Fridge, stove, payment system...',
    'Next step: Decides on its own',
]
add_bullet_list(slide, Inches(7.3), Inches(3.7), Inches(5.2), Inches(2.5),
                items_agent, size=15, color=TEXT_DARK)

# ════════════════════════════════════════════
# Agentic AI - How It Works
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_LIGHT)

add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
         "How Agentic AI Works", size=32, color=TEXT_DARK, bold=True)

steps = [
    ("1. Receive Goal", '"Create a monthly sales report"', ACCENT),
    ("2. Plan", "Query DB -> Analyze -> Visualize -> Write doc", ACCENT2),
    ("3. Execute Tools", "Run SQL queries, generate charts, draft text", ACCENT3),
    ("4. Evaluate", '"Chart is incomplete? Let me retry."', ORANGE),
    ("5. Iterate/Done", "Repeat 3-4 until resolved -> Report completion", N8N_COLOR),
]

for i, (step, desc, color) in enumerate(steps):
    y = Inches(1.3 + i * 1.1)
    add_shape_with_border(slide, Inches(1), y, Inches(11), Inches(0.9), BG_CARD, TEXT_LIGHT)

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), y, Inches(0.1), Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    add_text(slide, Inches(1.4), y + Inches(0.15), Inches(2.5), Inches(0.5),
             step, size=20, color=color, bold=True)
    add_text(slide, Inches(4.2), y + Inches(0.15), Inches(7), Inches(0.5),
             desc, size=17, color=TEXT_DARK)

add_text(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.5),
         "Key: Steps 2-5 are performed autonomously, without human intervention",
         size=16, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# Agentic AI - Core Traits
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_LIGHT)

add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
         "Six Core Characteristics of Agentic AI", size=32, color=TEXT_DARK, bold=True)

traits = [
    ("Autonomy", "Proceeds toward a goal without\nstep-by-step instructions", '"Create a report" -- handles everything', ACCENT),
    ("Multi-step\nReasoning", "Formulates and executes\ncomplex plans", "DB -> analysis -> visualization -> doc", ACCENT2),
    ("Tool Use", "Accesses APIs, databases,\nfile systems, the web", "Sends Slack messages, runs SQL queries", ACCENT3),
    ("Memory", "Maintains long-term and\nshort-term context", "Remembers yesterday's work today", YELLOW),
    ("Adaptability", "Revises plans when\nthings fail", "API error -> tries alternative path", ORANGE),
    ("Goal\nOrientation", "Repeats perceive-reason-act\nloop", "Keeps iterating until result is satisfactory", N8N_COLOR),
]

for i, (name, desc, example, color) in enumerate(traits):
    row = i // 3
    col = i % 3
    x = Inches(0.5 + col * 4.2)
    y = Inches(1.3 + row * 2.8)

    add_shape_with_border(slide, x, y, Inches(3.8), Inches(2.4), BG_CARD, TEXT_LIGHT)

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(3.8), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    add_text(slide, x + Inches(0.2), y + Inches(0.2), Inches(3.4), Inches(0.5),
             name, size=20, color=color, bold=True)
    add_text(slide, x + Inches(0.2), y + Inches(0.8), Inches(3.4), Inches(0.5),
             desc, size=14, color=TEXT_DARK)
    add_text(slide, x + Inches(0.2), y + Inches(1.5), Inches(3.4), Inches(0.6),
             f"e.g., {example}", size=12, color=TEXT_MID)

# ════════════════════════════════════════════
# AI Evolution Stages
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_LIGHT)

add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
         "The Evolution of AI", size=32, color=TEXT_DARK, bold=True)

generations = [
    ("Gen 1", "~2022", "Rule-based Chatbots", "FAQ matching", TEXT_LIGHT),
    ("Gen 2", "2023", "Generative AI", "Ask a question, get an answer (ChatGPT)", ACCENT2),
    ("Gen 3", "2024", "Copilots", "Suggests alongside you (Copilot)", ACCENT),
    ("Gen 4", "2025", "Single Agents", "Works independently (Claude Code)", ACCENT3),
    ("Gen 5", "2026", "Multi-Agent Systems", "Collaborates as a team (Agent Teams)", N8N_COLOR),
]

for i, (gen, year, name, desc, color) in enumerate(generations):
    y = Inches(1.3 + i * 1.15)
    bar_width = Inches(2 + i * 2)
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), y + Inches(0.15), bar_width, Inches(0.6))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    add_text(slide, Inches(0.5), y, Inches(1.2), Inches(0.5),
             gen, size=18, color=color, bold=True)
    add_text(slide, Inches(1.8), y, Inches(1.5), Inches(0.5),
             year, size=16, color=TEXT_MID)
    add_text(slide, Inches(4.2), y + Inches(0.15), bar_width, Inches(0.5),
             f"{name} -- {desc}", size=14, color=WHITE, bold=True)

add_text(slide, Inches(0.5), Inches(6.8), Inches(12), Inches(0.4),
         "We are at the Gen 4 -> Gen 5 transition right now",
         size=16, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# Single vs Multi-Agent
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_LIGHT)

add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
         "Single Agent vs Multi-Agent", size=32, color=TEXT_DARK, bold=True)

# Single
add_shape_with_border(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(3.5), BG_CARD, TEXT_LIGHT)
add_text(slide, Inches(0.8), Inches(1.6), Inches(5.2), Inches(0.5),
         "Single Agent", size=22, color=TEXT_MID, bold=True)
single_items = [
    "One AI handles everything",
    "Hits context window limits",
    "Versatile but lacks specialization",
    "",
    "DevOps case study:",
    "1.7% of recommendations were actionable",
]
add_bullet_list(slide, Inches(0.8), Inches(2.3), Inches(5.2), Inches(2.5),
                single_items, size=15, color=TEXT_MID)

# Multi
add_shape_with_border(slide, Inches(7), Inches(1.5), Inches(5.8), Inches(3.5), BG_CARD, ACCENT)
add_text(slide, Inches(7.3), Inches(1.6), Inches(5.2), Inches(0.5),
         "Multi-Agent System", size=22, color=ACCENT, bold=True)
multi_items = [
    "Specialized agents with distinct roles",
    "Parallel execution for speed",
    "Orchestrator coordinates the team",
    "",
    "DevOps case study:",
    "100% of recommendations were actionable",
]
add_bullet_list(slide, Inches(7.3), Inches(2.3), Inches(5.2), Inches(2.5),
                multi_items, size=15, color=TEXT_DARK)

# Bottom stats
stats_items = [
    "Gartner: Multi-agent system inquiries surged 1,445% (2024 Q1 -> 2025 Q2)",
    "Market size: $7.8B -> $52B+ (projected by 2030)",
    "Agentic AI success depends not on the model, but on orchestration",
]
add_bullet_list(slide, Inches(0.5), Inches(5.5), Inches(12), Inches(1.5),
                stats_items, size=15, color=ACCENT)

# ════════════════════════════════════════════
# Agent Framework Comparison
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_LIGHT)

add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
         "Major Agent Framework Comparison", size=32, color=TEXT_DARK, bold=True)

add_text(slide, Inches(0.5), Inches(1.1), Inches(12), Inches(0.4),
         "Tools for building agentic AI -- where does n8n fit?",
         size=16, color=TEXT_MID)

frameworks = [
    ("LangGraph", "Code (Python), graph state machines", "Complex custom agents", "High", LANGGRAPH_COLOR, "LG"),
    ("CrewAI", "Code (Python), role-based teams", "Multi-agent prototyping", "Medium", CREWAI_COLOR, "CR"),
    ("OpenAI SDK", "Code (Python), minimalist", "OpenAI ecosystem", "Low", OPENAI_SDK_COLOR, "OA"),
    ("Claude SDK", "Code, MCP-native", "Claude-based agents", "Medium", CLAUDE_SDK_COLOR, "CL"),
    ("n8n", "Visual no-code/low-code", "Business + dev teams", "Low", N8N_COLOR, "n8n"),
]

for i, (name, approach, target, difficulty, color, icon_text) in enumerate(frameworks):
    y = Inches(1.6 + i * 0.85)
    add_shape_with_border(slide, Inches(0.5), y, Inches(12), Inches(0.75), BG_CARD, TEXT_LIGHT)

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), y, Inches(0.1), Inches(0.75))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    add_brand_circle(slide, Inches(0.75), y + Inches(0.1), Inches(0.55), color, icon_text, WHITE, font_size=10)

    add_text(slide, Inches(1.5), y + Inches(0.1), Inches(2), Inches(0.5),
             name, size=17, color=color, bold=True)
    add_text(slide, Inches(3.5), y + Inches(0.1), Inches(3.5), Inches(0.5),
             approach, size=13, color=TEXT_DARK)
    add_text(slide, Inches(7.2), y + Inches(0.1), Inches(2.5), Inches(0.5),
             target, size=13, color=TEXT_MID)
    add_text(slide, Inches(9.8), y + Inches(0.1), Inches(2.7), Inches(0.5),
             difficulty, size=14, color=color, bold=True)

# Why n8n fits our team
add_shape_with_border(slide, Inches(0.5), Inches(5.9), Inches(12), Inches(1.5), BG_CARD, N8N_COLOR)
add_text(slide, Inches(0.8), Inches(5.95), Inches(11), Inches(0.4),
         "Why n8n is a strong fit for automotive industry IT teams", size=16, color=N8N_COLOR, bold=True)
n8n_fit_items = [
    "Self-hosted to protect customer/dealer data  |  Non-developers can modify workflows  |  500+ integrations (Slack, email, CRM)  |  Self-hosted is free vs SaaS at hundreds per month",
]
add_bullet_list(slide, Inches(0.8), Inches(6.35), Inches(11.5), Inches(0.7),
                n8n_fit_items, size=13, color=TEXT_DARK)

# ════════════════════════════════════════════
# Adoption Reality
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_LIGHT)

add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
         "Agentic AI Adoption Reality", size=32, color=TEXT_DARK, bold=True)

# Left: 11% big number
add_shape_with_border(slide, Inches(0.5), Inches(1.2), Inches(4.5), Inches(3.0), BG_CARD, TEXT_LIGHT)
add_text(slide, Inches(0.5), Inches(1.4), Inches(4.5), Inches(1.2),
         "11%", size=72, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(0.5), Inches(2.6), Inches(4.5), Inches(0.5),
         "have reached production", size=20, color=TEXT_DARK, align=PP_ALIGN.CENTER)
add_text(slide, Inches(0.5), Inches(3.2), Inches(4.5), Inches(0.5),
         "(Deloitte 2026 survey)", size=14, color=TEXT_MID, align=PP_ALIGN.CENTER)
add_text(slide, Inches(0.5), Inches(3.7), Inches(4.5), Inches(0.4),
         "Why do 89% fail to reach production?", size=15, color=ORANGE, bold=True, align=PP_ALIGN.CENTER)

# Right: bar chart
adoption = [
    ("Exploring", "30%", 2.4, TEXT_LIGHT),
    ("Piloting", "38%", 3.0, ACCENT2),
    ("Deploy-ready", "14%", 1.1, ACCENT),
    ("In Production", "11%", 0.9, ORANGE),
]
for i, (label, pct, width_val, color) in enumerate(adoption):
    y = Inches(1.3 + i * 0.85)
    add_text(slide, Inches(5.3), y + Inches(0.05), Inches(1.8), Inches(0.4),
             label, size=15, color=TEXT_DARK)
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(7.2), y, Inches(width_val), Inches(0.55))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    add_text(slide, Inches(7.2) + Inches(width_val) + Inches(0.15), y + Inches(0.05),
             Inches(1), Inches(0.4), pct, size=16, color=color, bold=True)

# Bottom: 5 conditions
add_text(slide, Inches(0.5), Inches(4.5), Inches(12), Inches(0.4),
         "Five Conditions for Reaching Production", size=18, color=TEXT_DARK, bold=True)

conditions = [
    ("1", "Data Pipeline Maturity", ACCENT, True),
    ("2", "Governance & Security", ACCENT2, True),
    ("3", "Legacy System API Access", ACCENT3, True),
    ("4", "Org-wide AI Literacy", YELLOW, False),
    ("5", "Incremental Rollout", ORANGE, False),
]
for i, (num, cond_text, color, n8n_helps) in enumerate(conditions):
    x = Inches(0.5 + i * 2.5)
    y = Inches(5.1)
    add_shape_with_border(slide, x, y, Inches(2.3), Inches(1.1), BG_CARD, TEXT_LIGHT)
    add_text(slide, x, y + Inches(0.05), Inches(2.3), Inches(0.4),
             num, size=22, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.1), y + Inches(0.45), Inches(2.1), Inches(0.55),
             cond_text, size=12, color=TEXT_DARK, align=PP_ALIGN.CENTER)
    if n8n_helps:
        add_text(slide, x, y + Inches(0.9), Inches(2.3), Inches(0.3),
                 "n8n helps", size=10, color=N8N_COLOR, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, Inches(0.5), Inches(6.5), Inches(12), Inches(0.5),
         "n8n directly supports data connectivity, API integration, and incremental adoption -- key enablers for production",
         size=14, color=N8N_COLOR, bold=True, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# n8n Introduction
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_LIGHT)

add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
         "n8n: A Platform That Makes Agentic AI Real", size=32, color=N8N_COLOR, bold=True)

add_brand_circle(slide, Inches(10.8), Inches(0.15), Inches(0.9), N8N_COLOR, "n8n", WHITE, font_size=16)

info_items = [
    ("What", "Fair-code workflow automation + AI platform"),
    ("Philosophy", '"The flexibility of code + the speed of no-code"'),
    ("Integrations", "500+ apps and services"),
    ("AI", "LangChain-based 70+ dedicated AI nodes"),
    ("Deployment", "Self-hosted (free) or Cloud"),
    ("Scale", "$2.5B valuation, Nvidia investment"),
]

for i, (key, val) in enumerate(info_items):
    row = i // 2
    col = i % 2
    x = Inches(0.5 + col * 6.3)
    y = Inches(1.3 + row * 1.0)
    add_shape_with_border(slide, x, y, Inches(6.0), Inches(0.85), BG_CARD, TEXT_LIGHT)
    add_text(slide, x + Inches(0.3), y + Inches(0.12), Inches(1.5), Inches(0.5),
             key, size=16, color=N8N_COLOR, bold=True)
    add_text(slide, x + Inches(1.8), y + Inches(0.12), Inches(4), Inches(0.5),
             val, size=15, color=TEXT_DARK)

# Key numbers
numbers = [
    ("$180M", "Series C (Nvidia participation)"),
    ("$40M+", "ARR (10x YoY)"),
    ("3,000+", "Enterprise Customers"),
    ("75%", "Customers Using AI Tools"),
]
for i, (num, label) in enumerate(numbers):
    x = Inches(0.5 + i * 3.2)
    y = Inches(5.0)
    add_shape_with_border(slide, x, y, Inches(2.8), Inches(1.8), BG_CARD, N8N_COLOR)
    add_text(slide, x, y + Inches(0.3), Inches(2.8), Inches(0.6),
             num, size=28, color=N8N_COLOR, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x, y + Inches(1.0), Inches(2.8), Inches(0.5),
             label, size=13, color=TEXT_MID, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# n8n x Agentic AI Requirement Matching
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_LIGHT)

add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
         "n8n x Agentic AI: Requirement Matching", size=32, color=TEXT_DARK, bold=True)

matchings = [
    ("Autonomy", "Trigger-based auto-execution\n(Webhooks, schedules, events)", ACCENT),
    ("Multi-step\nReasoning", "AI Agent node\nPlan & Execute pattern", ACCENT2),
    ("Tool Use", "500+ integrations + MCP\nClient/Server", ACCENT3),
    ("Memory", "Simple / PostgreSQL\n/ Redis memory nodes", YELLOW),
    ("Adaptability", "Conditional branching,\nerror handling, retry logic", ORANGE),
    ("Goal\nOrientation", "Workflow = complete pipeline\nfrom start to goal", N8N_COLOR),
]

for i, (trait, how, color) in enumerate(matchings):
    row = i // 3
    col = i % 3
    x = Inches(0.5 + col * 4.2)
    y = Inches(1.3 + row * 2.8)

    add_shape_with_border(slide, x, y, Inches(3.8), Inches(2.4), BG_CARD, TEXT_LIGHT)

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(3.8), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    add_text(slide, x + Inches(0.2), y + Inches(0.3), Inches(3.4), Inches(0.5),
             trait, size=20, color=color, bold=True)
    add_text(slide, x + Inches(0.2), y + Inches(1.0), Inches(3.4), Inches(1.2),
             how, size=15, color=TEXT_DARK)

# ════════════════════════════════════════════
# n8n Position Map
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_LIGHT)

add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
         "n8n's Unique Position", size=32, color=TEXT_DARK, bold=True)

add_text(slide, Inches(0.5), Inches(1.1), Inches(12), Inches(0.5),
         "Agentic AI Tool Spectrum", size=18, color=TEXT_MID)

# Spectrum bar
spectrum = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(0.5), Inches(2.0), Inches(12), Inches(0.6))
spectrum.fill.solid()
spectrum.fill.fore_color.rgb = BG_CARD
spectrum.line.color.rgb = TEXT_LIGHT
spectrum.line.width = Pt(0.5)

add_text(slide, Inches(0.5), Inches(2.0), Inches(3), Inches(0.5),
         "<- For Code Experts", size=13, color=TEXT_MID)
add_text(slide, Inches(9.5), Inches(2.0), Inches(3), Inches(0.5),
         "For Business Teams ->", size=13, color=TEXT_MID, align=PP_ALIGN.RIGHT)

tools_spectrum = [
    ("LangGraph", Inches(1.2), ACCENT2, "High"),
    ("CrewAI", Inches(3.5), ACCENT3, "Medium"),
    ("Agent SDK", Inches(5.5), ACCENT, "Medium"),
    ("n8n", Inches(8.0), N8N_COLOR, "Low"),
    ("Zapier", Inches(10.5), TEXT_LIGHT, "Low"),
]

for name, x, color, difficulty in tools_spectrum:
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, Inches(2.1), Inches(0.4), Inches(0.4))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    add_text(slide, x - Inches(0.5), Inches(2.7), Inches(1.4), Inches(0.4),
             name, size=14, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x - Inches(0.5), Inches(3.1), Inches(1.4), Inches(0.3),
             difficulty, size=11, color=TEXT_MID, align=PP_ALIGN.CENTER)

# n8n sweet spot
add_text(slide, Inches(2), Inches(3.8), Inches(9), Inches(0.5),
         'n8n = "Code flexibility + No-code speed" = Sweet Spot',
         size=20, color=N8N_COLOR, bold=True, align=PP_ALIGN.CENTER)

# Comparison table
comparisons = [
    ("vs LangGraph", "Visual builder, faster to build", "Less suited for complex state machines"),
    ("vs CrewAI", "Self-hosted, data sovereignty", "Less sophisticated crew structures"),
    ("vs Zapier", "JS/Python code support, cost-efficient", "Fewer integrations (500 vs 7,000)"),
]

for i, (vs, pro, con) in enumerate(comparisons):
    y = Inches(4.6 + i * 0.75)
    add_text(slide, Inches(0.5), y, Inches(2.3), Inches(0.5),
             vs, size=15, color=ACCENT, bold=True)
    add_text(slide, Inches(3), y, Inches(4.5), Inches(0.5),
             f"+ {pro}", size=14, color=ACCENT3)
    add_text(slide, Inches(7.8), y, Inches(4.5), Inches(0.5),
             f"- {con}", size=14, color=ORANGE)

# ════════════════════════════════════════════
# n8n + MCP
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_LIGHT)

add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
         "n8n = The Agentic MCP Hub", size=32, color=TEXT_DARK, bold=True)

add_text(slide, Inches(0.5), Inches(1.1), Inches(12), Inches(0.5),
         "n8n supports MCP bidirectionally (since April 2025)", size=18, color=TEXT_MID)

# Diagram
# Left: MCP consumption
add_shape_with_border(slide, Inches(0.5), Inches(2.0), Inches(3.5), Inches(3), BG_CARD, ACCENT)
add_text(slide, Inches(0.5), Inches(2.1), Inches(3.5), Inches(0.5),
         "External MCP Servers", size=16, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
ext_servers = ["GitHub", "Slack", "PostgreSQL", "File System"]
add_bullet_list(slide, Inches(0.8), Inches(2.7), Inches(3), Inches(2),
                ext_servers, size=14, color=TEXT_MID)

# Center: n8n
add_shape(slide, Inches(4.5), Inches(2.0), Inches(4.5), Inches(3), ACCENT2)
add_text(slide, Inches(4.5), Inches(2.2), Inches(4.5), Inches(0.5),
         "n8n Agent", size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
n8n_features = [
    "MCP Client -> Consume tools",
    "AI Agent -> Orchestration",
    "MCP Server -> Expose tools",
]
add_bullet_list(slide, Inches(4.8), Inches(2.9), Inches(4), Inches(2),
                n8n_features, size=14, color=WHITE)

# Right: MCP provision
add_shape_with_border(slide, Inches(9.5), Inches(2.0), Inches(3.5), Inches(3), BG_CARD, N8N_COLOR)
add_text(slide, Inches(9.5), Inches(2.1), Inches(3.5), Inches(0.5),
         "External MCP Clients", size=16, color=N8N_COLOR, bold=True, align=PP_ALIGN.CENTER)
ext_clients = ["Claude Desktop", "Cursor", "VS Code", "Other Agents"]
add_bullet_list(slide, Inches(9.8), Inches(2.7), Inches(3), Inches(2),
                ext_clients, size=14, color=TEXT_MID)

# Arrows
add_text(slide, Inches(3.8), Inches(3.0), Inches(0.8), Inches(0.5),
         "->", size=30, color=ACCENT, align=PP_ALIGN.CENTER)
add_text(slide, Inches(8.8), Inches(3.0), Inches(0.8), Inches(0.5),
         "->", size=30, color=N8N_COLOR, align=PP_ALIGN.CENTER)

# Workflow examples
add_text(slide, Inches(0.5), Inches(5.3), Inches(12), Inches(0.5),
         "Real-World Workflow Examples", size=20, color=TEXT_DARK, bold=True)

examples = [
    ("Customer Inquiry Agent", "Email received -> AI classifies (test drive/service/quote) -> CRM lookup -> auto-process", ACCENT),
    ("Market Intelligence", "Weekly schedule -> Gather competitor pricing/promos -> Trend report -> Slack", ACCENT3),
    ("Predictive Maintenance", "Sensor data -> Anomaly analysis -> Proactive alert to customer -> Service booking", ORANGE),
]

for i, (name, flow, color) in enumerate(examples):
    y = Inches(5.9 + i * 0.5)
    add_text(slide, Inches(0.8), y, Inches(2.5), Inches(0.4),
             name, size=14, color=color, bold=True)
    add_text(slide, Inches(3.5), y, Inches(9), Inches(0.4),
             flow, size=13, color=TEXT_MID)

# ════════════════════════════════════════════
# Case Studies
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_LIGHT)

add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
         "Automotive AI + n8n Case Studies", size=32, color=TEXT_DARK, bold=True)

cases = [
    ("Mercedes-Benz", "MBUX AI agent (in-vehicle voice + gesture)", "Multimodal agent in production", ACCENT),
    ("BMW", "Automated dealer support + Neue Klasse AI", "Reduced CS processing time", ACCENT3),
    ("Volvo", "Predictive maintenance AI -- sensor analysis", "Increased preventive maintenance", ORANGE),
    ("Delivery Hero", "Operational automation workflows (n8n)", "200 hours saved per month", ACCENT2),
    ("BeGlobal", "AI-powered proposal generation (n8n)", "10x scaling, under 1 minute", N8N_COLOR),
]

for i, (company, usecase, result, color) in enumerate(cases):
    y = Inches(1.3 + i * 1.15)
    add_shape_with_border(slide, Inches(0.5), y, Inches(12), Inches(1.0), BG_CARD, TEXT_LIGHT)

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), y, Inches(0.12), Inches(1.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    add_text(slide, Inches(0.9), y + Inches(0.15), Inches(2.5), Inches(0.5),
             company, size=20, color=color, bold=True)
    add_text(slide, Inches(3.5), y + Inches(0.15), Inches(4), Inches(0.5),
             usecase, size=16, color=TEXT_DARK)
    add_text(slide, Inches(8), y + Inches(0.15), Inches(4.5), Inches(0.5),
             result, size=18, color=ACCENT3, bold=True)

# ════════════════════════════════════════════
# Honest Assessment
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_LIGHT)

add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
         "An Honest Assessment: Opportunities & Challenges", size=32, color=TEXT_DARK, bold=True)

# Opportunities
add_shape_with_border(slide, Inches(0.5), Inches(1.3), Inches(5.8), Inches(4.5), BG_CARD, ACCENT3)
add_text(slide, Inches(0.8), Inches(1.4), Inches(5.2), Inches(0.5),
         "Opportunities", size=22, color=ACCENT3, bold=True)
pros = [
    "No-code/low-code democratizes agent building",
    "Self-hosted ensures data sovereignty",
    "500+ integrations connect legacy systems",
    "Bidirectional MCP support -> ecosystem hub",
    "Visual builder enables rapid prototyping",
    "Nvidia investment, $2.5B valuation",
]
add_bullet_list(slide, Inches(0.8), Inches(2.1), Inches(5.2), Inches(3.5),
                pros, size=15, color=TEXT_DARK, spacing=Pt(12))

# Challenges
add_shape_with_border(slide, Inches(7), Inches(1.3), Inches(5.8), Inches(4.5), BG_CARD, ORANGE)
add_text(slide, Inches(7.3), Inches(1.4), Inches(5.2), Inches(0.5),
         "Challenges", size=22, color=ORANGE, bold=True)
cons = [
    "Complex agent logic favors code frameworks",
    "Enterprise governance & security validation needed",
    "LLM token cost management required",
    "The MCP efficiency debate continues",
    "Production stability needs further validation",
    "No guarantees of success (only 11% in production)",
]
add_bullet_list(slide, Inches(7.3), Inches(2.1), Inches(5.2), Inches(3.5),
                cons, size=15, color=TEXT_MID, spacing=Pt(12))

add_text(slide, Inches(0.5), Inches(6.2), Inches(12), Inches(0.5),
         "Success or failure aside -- the attempt itself is a learning opportunity",
         size=18, color=N8N_COLOR, bold=True, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# Internal System Integration: APIs Are the Key
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_LIGHT)

add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
         "Internal System Integration: APIs Are the Key", size=32, color=TEXT_DARK, bold=True)

add_text(slide, Inches(0.5), Inches(1.1), Inches(12), Inches(0.5),
         "For agentic AI to actually work, internal systems must be accessible via APIs",
         size=17, color=TEXT_MID)

# Diagram
# [User Request]
add_shape(slide, Inches(0.5), Inches(2.3), Inches(2.0), Inches(0.8), ACCENT2)
add_text(slide, Inches(0.5), Inches(2.4), Inches(2.0), Inches(0.6),
         "User Request", size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, Inches(2.5), Inches(2.45), Inches(0.5), Inches(0.5),
         "->", size=20, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

# [n8n Agent]
add_shape(slide, Inches(3.0), Inches(2.3), Inches(2.2), Inches(0.8), N8N_COLOR)
add_text(slide, Inches(3.0), Inches(2.4), Inches(2.2), Inches(0.6),
         "n8n Agent", size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, Inches(5.2), Inches(2.45), Inches(0.5), Inches(0.5),
         "->", size=20, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

# [API Gateway]
add_shape(slide, Inches(5.7), Inches(2.3), Inches(2.0), Inches(0.8), ACCENT)
add_text(slide, Inches(5.7), Inches(2.4), Inches(2.0), Inches(0.6),
         "API Gateway", size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, Inches(7.7), Inches(2.45), Inches(0.5), Inches(0.5),
         "->", size=20, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

# Internal systems
internal_systems = [
    ("DMS (Dealer Mgmt)", ACCENT3),
    ("CRM (Customer Mgmt)", ACCENT2),
    ("ERP (Inventory/Finance)", ORANGE),
    ("Service Booking System", YELLOW),
]
for i, (sys_name, sys_color) in enumerate(internal_systems):
    y = Inches(1.8 + i * 0.7)
    add_shape_with_border(slide, Inches(8.3), y, Inches(4.2), Inches(0.6), BG_CARD, TEXT_LIGHT)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.3), y, Inches(0.08), Inches(0.6))
    bar.fill.solid()
    bar.fill.fore_color.rgb = sys_color
    bar.line.fill.background()
    add_text(slide, Inches(8.6), y + Inches(0.05), Inches(3.8), Inches(0.4),
             sys_name, size=14, color=TEXT_DARK)

# Current state & required actions
add_shape_with_border(slide, Inches(0.5), Inches(4.5), Inches(5.8), Inches(2.0), BG_CARD, ORANGE)
add_text(slide, Inches(0.8), Inches(4.6), Inches(5.2), Inches(0.4),
         "Current State", size=18, color=ORANGE, bold=True)
current_items = [
    "Many internal systems offer no API or limited access",
    "Data trapped in silos",
    "Reliant on manual data extraction and entry",
]
add_bullet_list(slide, Inches(0.8), Inches(5.1), Inches(5.2), Inches(1.2),
                current_items, size=14, color=TEXT_MID)

add_shape_with_border(slide, Inches(7), Inches(4.5), Inches(5.8), Inches(2.0), BG_CARD, ACCENT3)
add_text(slide, Inches(7.3), Inches(4.6), Inches(5.2), Inches(0.4),
         "Required Actions", size=18, color=ACCENT3, bold=True)
action_items = [
    "Develop and expose APIs for core systems",
    "Standardize the API Gateway",
    "Design data flows (connect via n8n)",
]
add_bullet_list(slide, Inches(7.3), Inches(5.1), Inches(5.2), Inches(1.2),
                action_items, size=14, color=TEXT_DARK)

add_text(slide, Inches(0.5), Inches(6.8), Inches(12), Inches(0.5),
         "The success of AI agents = the maturity of your API infrastructure",
         size=20, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# Big Picture Summary
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_LIGHT)

add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
         "The Big Picture", size=32, color=TEXT_DARK, bold=True)

flow_steps = [
    ("Models got smarter", "LLM competition", TEXT_LIGHT),
    ("They can use tools now", "MCP, native AI", ACCENT),
    ("They write code themselves", "CLI, vibe coding", ACCENT2),
    ("They manage their own context", "Context engineering", ACCENT3),
    ("They work independently", "Agentic AI", ORANGE),
    ("They collaborate as teams", "Multi-agent systems", N8N_COLOR),
]

for i, (main_text, sub, color) in enumerate(flow_steps):
    y = Inches(1.2 + i * 0.95)

    add_shape_with_border(slide, Inches(2), y, Inches(9), Inches(0.8), BG_CARD, TEXT_LIGHT)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), y, Inches(0.1), Inches(0.8))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    add_text(slide, Inches(2.4), y + Inches(0.1), Inches(5), Inches(0.5),
             main_text, size=18, color=TEXT_DARK, bold=True)
    add_text(slide, Inches(7.5), y + Inches(0.1), Inches(3), Inches(0.5),
             sub, size=15, color=color)

    if i < 5:
        add_text(slide, Inches(6.3), y + Inches(0.6), Inches(0.5), Inches(0.4),
                 "v", size=16, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)

# Conclusion
add_text(slide, Inches(0.5), Inches(7.0), Inches(12), Inches(0.4),
         "The platform that visually combines and executes all of this = n8n",
         size=18, color=N8N_COLOR, bold=True, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# Next Steps
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_LIGHT)

add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
         "Proposed Next Steps", size=32, color=TEXT_DARK, bold=True)

next_steps = [
    ("Week 1", "Install n8n self-hosted\nExplore basic workflows", "Get familiar", ACCENT),
    ("Week 2", "Build a simple automation\nwith the AI Agent node", "Validate potential", ACCENT2),
    ("Week 3", "Convert one real process\ninto an agentic workflow", "Apply to real work", ACCENT3),
    ("Week 4", "Team review\nDecide on expansion", "Go / No-Go", N8N_COLOR),
]

for i, (week, action, goal, color) in enumerate(next_steps):
    x = Inches(0.5 + i * 3.2)
    y = Inches(1.5)

    add_shape_with_border(slide, x, y, Inches(2.8), Inches(4.0), BG_CARD, color)

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(2.8), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    add_text(slide, x, y + Inches(0.3), Inches(2.8), Inches(0.5),
             week, size=24, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.2), y + Inches(1.0), Inches(2.4), Inches(1.5),
             action, size=15, color=TEXT_DARK, align=PP_ALIGN.CENTER)
    add_text(slide, x, y + Inches(3.0), Inches(2.8), Inches(0.5),
             goal, size=16, color=TEXT_MID, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# Closing
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_LIGHT)

add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
         '"The success of agentic AI depends\nnot on the model, but on orchestration."',
         size=32, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)

# Divider line
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(3.5), Inches(5), Pt(2))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT
line.line.fill.background()

add_text(slide, Inches(1), Inches(4.0), Inches(11), Inches(0.8),
         "We don't know yet whether this will succeed or fail.\nBut if we never try, we'll never know.",
         size=22, color=TEXT_MID, align=PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
         "The age of agentic AI has already begun,\nand n8n opens the door at the lowest threshold.",
         size=20, color=N8N_COLOR, bold=True, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# Q&A
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_LIGHT)

add_text(slide, Inches(1), Inches(2.5), Inches(11), Inches(1),
         "Q & A", size=48, color=TEXT_DARK, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(4.0), Inches(11), Inches(0.5),
         "Thank You", size=24, color=ACCENT, align=PP_ALIGN.CENTER)

# ── Speaker Notes ──
speaker_notes = {
    0: "Introduce the topic after greeting. 'Today I'll summarize recent AI trends and share concrete ideas our team can put into practice.'",
    1: "Pose the question and pause for 2-3 seconds. After capturing attention, transition with 'A lot has actually changed in just the past year.'",
    2: "Terminology slide. 'This presentation uses a lot of technical terms. Let me cover the essentials first. Feel free to ask about anything unfamiliar.'",
    3: "Timeline. 'The key takeaway from this timeline is the sheer pace of change. All of this happened in roughly one year since late 2024.'",
    6: "MCP overview. 'Think of the USB-C analogy. Just like every device once needed a different charger, every AI needed custom integrations with each service. MCP unifies that into one standard.'",
    7: "MCP limitations. 'Interestingly, just as the standard solidified, major services began embedding their own native AI. It's like wireless charging becoming more convenient right as USB-C went universal.'",
}

total = len(prs.slides)
speaker_notes[total - 7] = "'There are no guarantees, but the attempt itself is valuable. In our industry -- customer support, predictive maintenance, market analysis -- there are clear areas worth exploring.'"
speaker_notes[total - 3] = "'The core of agentic AI is orchestration -- how well you combine tools and AI systems. n8n lowers that barrier more than any other platform.'"
speaker_notes[total - 1] = "'Thank you. If you have questions or ideas for applying this to our work, please share them freely.'"

for idx, note in speaker_notes.items():
    if 0 <= idx < len(prs.slides):
        notes_slide = prs.slides[idx].notes_slide
        notes_slide.notes_text_frame.text = note

# ── Save ──
output = "/Users/mung/git/ppt-maker/AI_Trends_2025-2026_Light.pptx"
prs.save(output)
print(f"Slides generated: {output}")
print(f"   Total slides: {len(prs.slides)}")
