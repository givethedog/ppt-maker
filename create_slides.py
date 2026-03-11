#!/usr/bin/env python3
"""AI 트렌드 2025-2026 발표 슬라이드 생성"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── 색상 팔레트 ──
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)      # 진한 남색
BG_MID = RGBColor(0x16, 0x21, 0x3E)       # 중간 남색
ACCENT = RGBColor(0x00, 0xD2, 0xFF)       # 시안
ACCENT2 = RGBColor(0x7C, 0x3A, 0xED)      # 보라
ACCENT3 = RGBColor(0x10, 0xB9, 0x81)      # 에메랄드
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xA0, 0xA0, 0xB0)
ORANGE = RGBColor(0xFF, 0x6B, 0x35)
YELLOW = RGBColor(0xFF, 0xD6, 0x00)
N8N_COLOR = RGBColor(0xEA, 0x4B, 0x71)    # n8n 브랜드

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_text(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name="Apple SD Gothic Neo"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_bullet_list(slide, left, top, width, height, items, size=16, color=WHITE, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Apple SD Gothic Neo"
        p.space_after = spacing
    return txBox


# ════════════════════════════════════════════
# 슬라이드 1: 타이틀
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_bg(slide, BG_DARK)

add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
         "AI 트렌드 2025-2026", size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(2.8), Inches(11), Inches(1),
         "에이전틱 AI 시대와 n8n의 가능성", size=28, color=ACCENT, align=PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.5),
         "IT팀 내부 공유  |  2026년 3월", size=18, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# 하단 장식 라인
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3), Inches(4.0), Inches(7), Pt(3))
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()

# ════════════════════════════════════════════
# 슬라이드 2: 오프닝 질문
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

add_text(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.5),
         '"지난 1년, AI 업계에서\n무슨 일이 있었을까요?"', size=40, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, Inches(2), Inches(4.5), Inches(9), Inches(1),
         "ChatGPT 이후 2년... 세상은 상상 이상으로 빠르게 바뀌고 있습니다",
         size=20, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# 슬라이드 3: 용어 설명
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
         "오늘 발표에서 나올 주요 용어", size=32, color=WHITE, bold=True)

add_text(slide, Inches(0.5), Inches(1.0), Inches(12), Inches(0.4),
         "AI/IT 용어가 많이 등장합니다. 핵심만 먼저 짚고 갈게요!",
         size=16, color=LIGHT_GRAY)

terms = [
    ("LLM", "Large Language Model", "ChatGPT 같은 대규모 언어 AI 모델", ACCENT),
    ("에이전트", "Agent", "스스로 판단하고 행동하는 AI 프로그램", ACCENT2),
    ("MCP", "Model Context Protocol", "AI가 외부 도구에 접속하는 표준 규격 (USB-C 같은 것)", ACCENT3),
    ("CLI", "Command Line Interface", "터미널/명령줄 — 마우스 대신 텍스트로 조작", YELLOW),
    ("토큰", "Token", "AI가 텍스트를 처리하는 단위 (≈ 한글 1글자 or 영어 4글자)", ORANGE),
    ("RAG", "Retrieval Augmented Generation", "AI가 외부 문서를 검색해서 답변에 활용하는 기술", N8N_COLOR),
    ("워크플로우", "Workflow", "업무 처리 흐름을 자동화한 것 (A→B→C 순서대로 실행)", ACCENT),
    ("오케스트레이션", "Orchestration", "여러 AI/도구를 지휘·조율하는 것 (지휘자 역할)", ACCENT2),
]

for i, (term_kr, term_en, desc, color) in enumerate(terms):
    row = i // 2
    col = i % 2
    x = Inches(0.5 + col * 6.3)
    y = Inches(1.6 + row * 1.35)

    add_shape(slide, x, y, Inches(6.0), Inches(1.2), BG_MID)

    # 색상 왼쪽 바
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.1), Inches(1.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    add_text(slide, x + Inches(0.3), y + Inches(0.05), Inches(2.5), Inches(0.45),
             f"{term_kr} ({term_en})", size=15, color=color, bold=True)
    add_text(slide, x + Inches(0.3), y + Inches(0.5), Inches(5.4), Inches(0.6),
             desc, size=13, color=LIGHT_GRAY)

add_text(slide, Inches(0.5), Inches(7.0), Inches(12), Inches(0.4),
         "💡 발표 중 모르는 용어가 나오면 언제든 질문해 주세요!",
         size=14, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# 슬라이드 4: 타임라인
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
         "최근 1년 타임라인", size=32, color=WHITE, bold=True)

timeline = [
    ("2024.11", "MCP 발표", ACCENT),
    ("2025.01", "DeepSeek R1 충격", ORANGE),
    ("2025.02", "바이브 코딩 등장", YELLOW),
    ("2025.03", "OpenAI Agents SDK", ACCENT2),
    ("2025.05", "Claude Code 출시", ACCENT),
    ("2025.09", "Notion AI Agents", ACCENT3),
    ("2025.10", "n8n $180M 투자", N8N_COLOR),
    ("2025.12", "MCP → Linux 재단", LIGHT_GRAY),
    ("2026.01", "오픈클로 25만 스타", ORANGE),
    ("2026.02", "Claude 4.6 출시", ACCENT),
]

for i, (date, event, color) in enumerate(timeline):
    row = i // 5
    col = i % 5
    x = Inches(0.4 + col * 2.5)
    y = Inches(1.5 + row * 2.8)

    # 원형 포인트
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.9), y, Inches(0.25), Inches(0.25))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()

    add_text(slide, x, y + Inches(0.4), Inches(2.2), Inches(0.4),
             date, size=13, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x, y + Inches(0.75), Inches(2.2), Inches(0.5),
             event, size=14, color=WHITE, align=PP_ALIGN.CENTER)

# 하단 메시지
add_text(slide, Inches(0.5), Inches(6.7), Inches(12), Inches(0.5),
         'AI: "대화 도구" → "일하는 에이전트" → "시스템을 운영하는 팀"으로 진화 중',
         size=16, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# 슬라이드 4: LLM 모델 경쟁
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
         "LLM 모델 경쟁: 5파전", size=32, color=WHITE, bold=True)

players = [
    ("OpenAI", "GPT-4.1, o3, o4-mini", "추론 특화, 1M 토큰", ACCENT3),
    ("Anthropic", "Claude 4.6", "코딩 1위, 30h+ 자율 에이전트", ACCENT),
    ("Google", "Gemini 2.5 Pro", "1M 토큰, Deep Think", YELLOW),
    ("오픈소스", "DeepSeek, Llama 4", "가격 파괴, 10M 토큰", ORANGE),
    ("챌린저", "Grok 3, Mistral", "인프라 + EU 데이터 주권", ACCENT2),
]

for i, (name, models, strength, color) in enumerate(players):
    y = Inches(1.3 + i * 1.15)
    card = add_shape(slide, Inches(0.5), y, Inches(12), Inches(1.0), BG_MID)

    # 왼쪽 색상 바
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), y, Inches(0.12), Inches(1.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    add_text(slide, Inches(0.9), y + Inches(0.1), Inches(2), Inches(0.5),
             name, size=20, color=color, bold=True)
    add_text(slide, Inches(3.2), y + Inches(0.1), Inches(3.5), Inches(0.5),
             models, size=16, color=WHITE)
    add_text(slide, Inches(7), y + Inches(0.1), Inches(5.5), Inches(0.5),
             strength, size=16, color=LIGHT_GRAY)

# 하단
add_text(slide, Inches(0.5), Inches(6.8), Inches(12), Inches(0.5),
         "경쟁의 핵심: 벤치마크 → 에이전트 능력 → 팀 조율 능력으로 전환",
         size=15, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# 슬라이드 5: 경쟁 핵심 트렌드
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
         "모델 경쟁의 핵심 축 변화", size=32, color=WHITE, bold=True)

# 화살표 흐름
phases = [
    ("2024", "벤치마크 점수 경쟁", LIGHT_GRAY),
    ("2025", "얼마나 오래 자율 작업 가능?", ACCENT),
    ("2026", "에이전트 팀을 얼마나 잘 조율?", ACCENT2),
]
for i, (year, desc, color) in enumerate(phases):
    y = Inches(1.5 + i * 1.6)
    add_text(slide, Inches(1.5), y, Inches(2), Inches(0.6),
             year, size=36, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(4), y + Inches(0.05), Inches(7), Inches(0.6),
             desc, size=22, color=WHITE)
    if i < 2:
        add_text(slide, Inches(2.2), y + Inches(0.8), Inches(1), Inches(0.5),
                 "↓", size=30, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

trends = [
    "✦ 모든 모델이 추론(thinking) 모드 탑재 → 기본 기능화",
    "✦ DeepSeek R1 이후 전반적 가격 하락 → 비용 장벽 소멸",
    "✦ 오픈소스가 클로즈드 모델 성능에 급격히 근접",
]
add_bullet_list(slide, Inches(1), Inches(5.8), Inches(11), Inches(1.5),
                trends, size=15, color=LIGHT_GRAY)

# ════════════════════════════════════════════
# 슬라이드 6: MCP 개요
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
         "MCP: AI의 USB-C", size=32, color=WHITE, bold=True)

add_text(slide, Inches(0.5), Inches(1.2), Inches(12), Inches(0.6),
         "Model Context Protocol — AI가 외부 도구에 접속하는 범용 표준",
         size=18, color=LIGHT_GRAY)

# Before
add_shape(slide, Inches(0.5), Inches(2.2), Inches(5.8), Inches(3.5), BG_MID)
add_text(slide, Inches(0.8), Inches(2.3), Inches(5), Inches(0.5),
         "Before MCP", size=20, color=ORANGE, bold=True)
before_items = [
    "AI 모델 A ─── 커스텀 코드 ─── GitHub",
    "AI 모델 A ─── 커스텀 코드 ─── Slack",
    "AI 모델 B ─── 커스텀 코드 ─── GitHub",
    "...",
    "M × N 개의 통합 코드 필요 😱",
]
add_bullet_list(slide, Inches(0.8), Inches(2.9), Inches(5), Inches(2.5),
                before_items, size=14, color=LIGHT_GRAY)

# After
add_shape(slide, Inches(7), Inches(2.2), Inches(5.8), Inches(3.5), BG_MID)
add_text(slide, Inches(7.3), Inches(2.3), Inches(5), Inches(0.5),
         "After MCP", size=20, color=ACCENT3, bold=True)
after_items = [
    "AI 모델 A ─┐",
    "AI 모델 B ─┼─ MCP ─┬─ GitHub Server",
    "AI 모델 C ─┘       ├─ Slack Server",
    "                    └─ DB Server",
    "M + N 개의 구현으로 충분! ✨",
]
add_bullet_list(slide, Inches(7.3), Inches(2.9), Inches(5), Inches(2.5),
                after_items, size=14, color=LIGHT_GRAY)

# 숫자들
stats = [
    ("5,800+", "MCP 서버"),
    ("300+", "클라이언트"),
    ("9,700만+", "월간 다운로드"),
]
for i, (num, label) in enumerate(stats):
    x = Inches(1.5 + i * 4)
    add_text(slide, x, Inches(6.0), Inches(3), Inches(0.5),
             num, size=28, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x, Inches(6.5), Inches(3), Inches(0.4),
             label, size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# 슬라이드 7: MCP의 한계
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
         "그런데... MCP는 이미 사양의 기미?", size=32, color=ORANGE, bold=True)

add_text(slide, Inches(0.5), Inches(1.2), Inches(12), Inches(0.6),
         "주요 서비스들이 자체 AI 에이전트를 내장하기 시작",
         size=18, color=LIGHT_GRAY)

services = [
    ("Notion", "3.0 AI Agents → 3.3 Custom Agents", "네이티브 주력, MCP도 병행"),
    ("Linear", "GitHub Copilot Agent 네이티브 통합", "MCP 없이 직접 연동"),
    ("GitHub", "Copilot CLI + Agent Mode", "자체 CLI가 MCP보다 효율적"),
]

for i, (name, feature, note) in enumerate(services):
    y = Inches(2.2 + i * 1.3)
    add_shape(slide, Inches(0.5), y, Inches(12), Inches(1.1), BG_MID)
    add_text(slide, Inches(0.9), y + Inches(0.15), Inches(2), Inches(0.4),
             name, size=20, color=ACCENT, bold=True)
    add_text(slide, Inches(3.2), y + Inches(0.15), Inches(5), Inches(0.4),
             feature, size=16, color=WHITE)
    add_text(slide, Inches(8.5), y + Inches(0.15), Inches(4), Inches(0.4),
             note, size=14, color=LIGHT_GRAY)

add_text(slide, Inches(0.5), Inches(6.0), Inches(12), Inches(0.4),
         "→ 주요 서비스들이 자체 AI를 내장하면서 MCP 의존도가 줄어드는 추세",
         size=15, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# 슬라이드 9: CLI vs MCP 효율성
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
         "CLI vs MCP: 효율성 논쟁", size=32, color=WHITE, bold=True)

add_text(slide, Inches(0.5), Inches(1.2), Inches(12), Inches(0.5),
         "CLI 도구가 MCP보다 AI 에이전트에게 더 효율적일 수 있다",
         size=18, color=LIGHT_GRAY)

# 비교 카드 - CLI
add_shape(slide, Inches(0.5), Inches(2.2), Inches(5.8), Inches(3.0), BG_MID)
add_text(slide, Inches(0.8), Inches(2.3), Inches(5.2), Inches(0.5),
         "CLI 방식 (gh, kubectl 등)", size=20, color=ACCENT3, bold=True)
cli_items = [
    "토큰 비용: ~200 토큰",
    "LLM이 이미 학습 데이터에서 잘 알고 있음",
    "벤치마크 33% 우위",
    "예: gh pr list → 바로 실행",
]
add_bullet_list(slide, Inches(0.8), Inches(3.0), Inches(5.2), Inches(2.0),
                cli_items, size=15, color=WHITE)

# 비교 카드 - MCP
add_shape(slide, Inches(7), Inches(2.2), Inches(5.8), Inches(3.0), BG_MID)
add_text(slide, Inches(7.3), Inches(2.3), Inches(5.2), Inches(0.5),
         "MCP 방식", size=20, color=ORANGE, bold=True)
mcp_items = [
    "토큰 비용: ~55,000 토큰 (GitHub 93개 도구)",
    "런타임에 처음 만나는 스키마",
    "컨텍스트 소모 과다",
    "범용성은 높지만 비효율적",
]
add_bullet_list(slide, Inches(7.3), Inches(3.0), Inches(5.2), Inches(2.0),
                mcp_items, size=15, color=LIGHT_GRAY)

# 결론
add_shape(slide, Inches(1.5), Inches(5.8), Inches(10), Inches(1.2), BG_MID)
add_text(slide, Inches(1.8), Inches(5.9), Inches(9.4), Inches(1.0),
         "MCP는 사라지지 않는다 (이미 표준)\n"
         "하지만 \"만능 열쇠\"에서 \"여러 방법 중 하나\"로 변화\n"
         "진짜 승자: 이 모든 것을 조합하는 오케스트레이션 레이어 ⭐",
         size=16, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# 슬라이드 10: CLI 회귀 & 바이브 코딩
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
         "CLI로의 회귀 & 바이브 코딩", size=32, color=WHITE, bold=True)

# 진화 단계
stages = [
    ("2023", "자동완성", "Copilot, Tabnine", LIGHT_GRAY),
    ("2024", "대화형 IDE", "Cursor, Windsurf", ACCENT2),
    ("2025", "CLI 에이전트", "Claude Code, Aider", ACCENT),
    ("2026", "팀 에이전트", "Claude Teams, n8n", N8N_COLOR),
]
for i, (year, stage, tools, color) in enumerate(stages):
    x = Inches(0.5 + i * 3.2)
    add_shape(slide, x, Inches(1.5), Inches(2.8), Inches(1.8), BG_MID)
    add_text(slide, x, Inches(1.6), Inches(2.8), Inches(0.4),
             year, size=24, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x, Inches(2.1), Inches(2.8), Inches(0.4),
             stage, size=18, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, x, Inches(2.6), Inches(2.8), Inches(0.4),
             tools, size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    if i < 3:
        add_text(slide, x + Inches(2.8), Inches(2.0), Inches(0.4), Inches(0.5),
                 "→", size=24, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Claude Code
add_text(slide, Inches(0.5), Inches(3.8), Inches(6), Inches(0.5),
         "Claude Code: 8개월 만에 선호도 1위 (46%)", size=18, color=ACCENT, bold=True)

# 바이브 코딩
add_shape(slide, Inches(0.5), Inches(4.5), Inches(6), Inches(2.5), BG_MID)
add_text(slide, Inches(0.8), Inches(4.6), Inches(5.5), Inches(0.5),
         '바이브 코딩 (Karpathy, 2025.02)', size=18, color=YELLOW, bold=True)
add_text(slide, Inches(0.8), Inches(5.1), Inches(5.5), Inches(0.8),
         '"코드가 존재한다는 사실조차 잊어라"\n→ 오픈클로: 25만 스타, 코드 안 읽고 배포\n→ 보안 취약점 2.74배 ↑',
         size=14, color=LIGHT_GRAY)

add_shape(slide, Inches(7), Inches(4.5), Inches(5.8), Inches(2.5), BG_MID)
add_text(slide, Inches(7.3), Inches(4.6), Inches(5.2), Inches(0.5),
         '에이전틱 엔지니어링 (2026)', size=18, color=ACCENT3, bold=True)
add_text(slide, Inches(7.3), Inches(5.1), Inches(5.2), Inches(0.8),
         '"바이브 코딩은 passé" — Karpathy\n→ AI가 코드를 쓰되, 사람이 감독·조율\n→ 품질과 보안도 중요',
         size=14, color=LIGHT_GRAY)

add_text(slide, Inches(6.2), Inches(5.3), Inches(0.6), Inches(0.5),
         "→", size=28, color=ACCENT, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# 슬라이드 9: 컨텍스트 엔지니어링
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
         "프롬프트 → 컨텍스트 엔지니어링", size=32, color=WHITE, bold=True)

add_text(slide, Inches(0.5), Inches(1.1), Inches(12), Inches(0.5),
         '"AI의 새로운 핵심 역량은 프롬프팅이 아니라 컨텍스트 엔지니어링" — Tobi Lutke (Shopify CEO)',
         size=16, color=LIGHT_GRAY)

# 비교 테이블
headers = ["", "프롬프트 엔지니어링", "컨텍스트 엔지니어링"]
rows = [
    ["범위", "하나의 텍스트 문자열", "전체 정보 아키텍처"],
    ["시점", "호출 시 1회", "지속적, 여러 턴에 걸쳐"],
    ["초점", "무엇을 물어볼까", "AI가 이미 무엇을 알고 있을까"],
    ["결과", "더 나은 문구", "더 나은 시스템 설계"],
]

for i, header in enumerate(headers):
    x = Inches(0.5 + i * 4.2)
    add_shape(slide, x, Inches(1.8), Inches(4.0), Inches(0.6), ACCENT2 if i > 0 else BG_MID)
    add_text(slide, x + Inches(0.2), Inches(1.85), Inches(3.6), Inches(0.5),
             header, size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

for ri, row in enumerate(rows):
    for ci, cell in enumerate(row):
        x = Inches(0.5 + ci * 4.2)
        y = Inches(2.5 + ri * 0.65)
        bg = BG_MID if ri % 2 == 0 else BG_DARK
        add_shape(slide, x, y, Inches(4.0), Inches(0.6), bg)
        color = ACCENT if ci == 2 else (LIGHT_GRAY if ci == 1 else WHITE)
        add_text(slide, x + Inches(0.2), y + Inches(0.05), Inches(3.6), Inches(0.5),
                 cell, size=14, color=color, align=PP_ALIGN.CENTER)

add_text(slide, Inches(0.5), Inches(5.2), Inches(12), Inches(0.4),
         "핵심: 프롬프트 = '무엇을 물어볼까'  vs  컨텍스트 = 'AI가 이미 무엇을 알고 있을까'",
         size=16, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# 컨텍스트 엔지니어링 - 7요소 (분리된 슬라이드)
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
         "컨텍스트의 7요소", size=32, color=WHITE, bold=True)

add_text(slide, Inches(0.5), Inches(1.1), Inches(12), Inches(0.5),
         "AI에게 제공하는 정보를 7가지 층으로 설계합니다",
         size=18, color=LIGHT_GRAY)

elements_detail = [
    ("시스템 프롬프트", "AI의 역할과 규칙 정의", "\"너는 고객 응대 전문가야\"", ACCENT),
    ("사용자 프롬프트", "현재 과제/질문", "\"이 이메일에 답장 써줘\"", ACCENT2),
    ("대화 이력", "지금까지의 맥락", "이전 대화 내용 참조", ACCENT3),
    ("장기 메모리", "세션을 넘는 기억", "지난주 작업 결과 기억", N8N_COLOR),
    ("검색된 정보 (RAG)", "필요한 문서·데이터", "사내 매뉴얼 검색 후 참조", ORANGE),
    ("사용 가능한 도구", "API, MCP, 함수", "Slack 전송, DB 조회 가능", ACCENT),
    ("출력 형식", "구조화된 응답 정의", "JSON, 표, 마크다운 등", YELLOW),
]

for i, (name, desc, example, color) in enumerate(elements_detail):
    row = i // 2
    col = i % 2
    x = Inches(0.5 + col * 6.3)
    y = Inches(1.8 + row * 1.3)

    add_shape(slide, x, y, Inches(6.0), Inches(1.15), BG_MID)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(0.1), Inches(1.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    add_text(slide, x + Inches(0.3), y + Inches(0.05), Inches(2.5), Inches(0.4),
             name, size=16, color=color, bold=True)
    add_text(slide, x + Inches(3.0), y + Inches(0.05), Inches(2.8), Inches(0.4),
             desc, size=13, color=WHITE)
    add_text(slide, x + Inches(0.3), y + Inches(0.55), Inches(5.5), Inches(0.4),
             f"예: {example}", size=12, color=LIGHT_GRAY)

add_text(slide, Inches(0.5), Inches(7.0), Inches(12), Inches(0.4),
         "핵심: 컨텍스트 엔지니어링 = 무엇을, 언제, 어떤 형식으로, 얼마나 가져올지 설계하는 상위 학문",
         size=14, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# 브릿지 슬라이드: 지금까지 → 에이전틱 AI
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

add_text(slide, Inches(1), Inches(1.0), Inches(11), Inches(0.8),
         "여기까지 정리하면...", size=36, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

bridge_items = [
    ("모델은 충분히 똑똑해졌고", LIGHT_GRAY),
    ("도구 연결 표준(MCP)도 만들어졌고", ACCENT),
    ("코드도 직접 쓰게 됐고", ACCENT2),
    ("컨텍스트 관리 방법도 진화했다", ACCENT3),
]

for i, (text, color) in enumerate(bridge_items):
    y = Inches(2.3 + i * 0.9)
    add_text(slide, Inches(2), y, Inches(9), Inches(0.6),
             f"✓  {text}", size=22, color=color)

add_text(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.8),
         "그렇다면 이 모든 것이 합쳐지면?", size=28, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(6.3), Inches(11), Inches(0.8),
         "→  에이전틱 AI", size=40, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# 에이전틱 AI - 정의
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
         "에이전틱 AI: 그래서 뭔데?", size=32, color=WHITE, bold=True)

add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(0.8),
         "생성형 AI는 질문에 \"답\"한다.\n에이전틱 AI는 목표를 \"달성\"한다.",
         size=28, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

# 식당 비유
add_shape(slide, Inches(0.5), Inches(3.0), Inches(5.8), Inches(3.8), BG_MID)
add_text(slide, Inches(0.8), Inches(3.1), Inches(5.2), Inches(0.5),
         "🍳 생성형 AI = 레시피 알려주는 셰프", size=18, color=LIGHT_GRAY, bold=True)
items_gen = [
    '입력: "파스타 레시피 알려줘"',
    '출력: 텍스트 (레시피)',
    '도구: 없음',
    '다음 단계: 사람이 지시',
]
add_bullet_list(slide, Inches(0.8), Inches(3.7), Inches(5.2), Inches(2.5),
                items_gen, size=15, color=LIGHT_GRAY)

add_shape(slide, Inches(7), Inches(3.0), Inches(5.8), Inches(3.8), BG_MID)
add_text(slide, Inches(7.3), Inches(3.1), Inches(5.2), Inches(0.5),
         "🤖 에이전틱 AI = 요리+서빙까지 하는 셰프", size=18, color=ACCENT, bold=True)
items_agent = [
    '입력: "저녁 파티 준비해줘"',
    '출력: 행동 (장보기→요리→세팅→서빙)',
    '도구: 냉장고, 가스레인지, 카드결제...',
    '다음 단계: 스스로 결정',
]
add_bullet_list(slide, Inches(7.3), Inches(3.7), Inches(5.2), Inches(2.5),
                items_agent, size=15, color=WHITE)

# ════════════════════════════════════════════
# 슬라이드 11: 에이전틱 AI - 작동 원리
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
         "에이전틱 AI 작동 원리", size=32, color=WHITE, bold=True)

steps = [
    ("① 목표 수신", '"월간 매출 리포트 만들어줘"', ACCENT),
    ("② 계획 수립", "DB 조회 → 분석 → 시각화 → 문서 작성", ACCENT2),
    ("③ 도구 실행", "SQL 쿼리, 차트 생성, 문서 작성", ACCENT3),
    ("④ 결과 평가", '"차트가 불완전? 다시 해볼게"', ORANGE),
    ("⑤ 반복/완료", "문제 해결까지 ③-④ 반복 → 완료 보고", N8N_COLOR),
]

for i, (step, desc, color) in enumerate(steps):
    y = Inches(1.3 + i * 1.1)
    add_shape(slide, Inches(1), y, Inches(11), Inches(0.9), BG_MID)

    # 색상 바
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), y, Inches(0.1), Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    add_text(slide, Inches(1.4), y + Inches(0.15), Inches(2.5), Inches(0.5),
             step, size=20, color=color, bold=True)
    add_text(slide, Inches(4.2), y + Inches(0.15), Inches(7), Inches(0.5),
             desc, size=17, color=WHITE)

    if i < 4 and i >= 2:
        # 반복 화살표
        pass

add_text(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.5),
         "핵심: ②~⑤를 사람 개입 없이 자율적으로 수행하는 것이 에이전틱 AI",
         size=16, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# 슬라이드 12: 에이전틱 AI - 핵심 특성
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
         "에이전틱 AI 핵심 특성 6가지", size=32, color=WHITE, bold=True)

traits = [
    ("자율성", "단계별 지시 없이 목표를 향해 진행", '"리포트 만들어줘" 한 마디로 끝까지', ACCENT),
    ("다단계 추론", "복잡한 계획을 세우고 실행", "DB→분석→시각화→문서 파이프라인", ACCENT2),
    ("도구 사용", "API, DB, 파일시스템, 웹 접근", "Slack 메시지 전송, SQL 쿼리 실행", ACCENT3),
    ("메모리", "장기·단기 컨텍스트 유지", "어제 작업 내용을 오늘도 기억", YELLOW),
    ("적응력", "실패 시 계획 수정", "API 오류 → 대안 경로 시도", ORANGE),
    ("목표 지향", "인식→추론→행동 루프 반복", "결과가 만족스러울 때까지 반복", N8N_COLOR),
]

for i, (name, desc, example, color) in enumerate(traits):
    row = i // 3
    col = i % 3
    x = Inches(0.5 + col * 4.2)
    y = Inches(1.3 + row * 2.8)

    add_shape(slide, x, y, Inches(3.8), Inches(2.4), BG_MID)

    # 색상 상단 바
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(3.8), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    add_text(slide, x + Inches(0.2), y + Inches(0.2), Inches(3.4), Inches(0.5),
             name, size=20, color=color, bold=True)
    add_text(slide, x + Inches(0.2), y + Inches(0.8), Inches(3.4), Inches(0.5),
             desc, size=14, color=WHITE)
    add_text(slide, x + Inches(0.2), y + Inches(1.5), Inches(3.4), Inches(0.6),
             f"예: {example}", size=12, color=LIGHT_GRAY)

# ════════════════════════════════════════════
# 슬라이드 13: AI 진화 단계
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
         "AI의 진화 단계", size=32, color=WHITE, bold=True)

generations = [
    ("1세대", "~2022", "규칙 기반 챗봇", "FAQ 매칭", LIGHT_GRAY),
    ("2세대", "2023", "생성형 AI", "질문하면 답변 (ChatGPT)", ACCENT2),
    ("3세대", "2024", "코파일럿", "옆에서 제안 (Copilot)", ACCENT),
    ("4세대", "2025", "단일 에이전트", "혼자서 일 처리 (Claude Code)", ACCENT3),
    ("5세대", "2026", "멀티 에이전트", "팀으로 협업 (Agent Teams)", N8N_COLOR),
]

for i, (gen, year, name, desc, color) in enumerate(generations):
    y = Inches(1.3 + i * 1.15)
    # 진행 바
    bar_width = Inches(2 + i * 2)
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), y + Inches(0.15), bar_width, Inches(0.6))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    add_text(slide, Inches(0.5), y, Inches(1.2), Inches(0.5),
             gen, size=18, color=color, bold=True)
    add_text(slide, Inches(1.8), y, Inches(1.5), Inches(0.5),
             year, size=16, color=LIGHT_GRAY)
    add_text(slide, Inches(4.2), y + Inches(0.15), bar_width, Inches(0.5),
             f"{name} — {desc}", size=14, color=BG_DARK, bold=True)

add_text(slide, Inches(0.5), Inches(6.8), Inches(12), Inches(0.4),
         "우리는 지금 4세대 → 5세대 전환기에 있습니다",
         size=16, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# 슬라이드 14: 싱글 vs 멀티 에이전트
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
         "싱글 에이전트 vs 멀티 에이전트", size=32, color=WHITE, bold=True)

# 싱글
add_shape(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(3.5), BG_MID)
add_text(slide, Inches(0.8), Inches(1.6), Inches(5.2), Inches(0.5),
         "싱글 에이전트", size=22, color=LIGHT_GRAY, bold=True)
single_items = [
    "하나의 AI가 모든 것을 처리",
    "컨텍스트 윈도우 한계에 부딪힘",
    "범용적이지만 전문성 부족",
    "",
    "DevOps 사례:",
    "실행 가능 권고안 1.7% 😢",
]
add_bullet_list(slide, Inches(0.8), Inches(2.3), Inches(5.2), Inches(2.5),
                single_items, size=15, color=LIGHT_GRAY)

# 멀티
add_shape(slide, Inches(7), Inches(1.5), Inches(5.8), Inches(3.5), BG_MID)
add_text(slide, Inches(7.3), Inches(1.6), Inches(5.2), Inches(0.5),
         "멀티 에이전트 시스템", size=22, color=ACCENT, bold=True)
multi_items = [
    "전문 에이전트들이 역할 분담",
    "병렬 실행으로 속도 향상",
    "오케스트레이터가 조율",
    "",
    "DevOps 사례:",
    "실행 가능 권고안 100% 🎯",
]
add_bullet_list(slide, Inches(7.3), Inches(2.3), Inches(5.2), Inches(2.5),
                multi_items, size=15, color=WHITE)

# 하단 통계
stats_items = [
    "Gartner: 멀티 에이전트 시스템 문의 1,445% 급증 (2024 Q1 → 2025 Q2)",
    "시장 규모: $7.8B → $52B+ (2030 예측)",
    "에이전틱 AI 성패 = 모델이 아니라 오케스트레이션(조율)에 달려 있다",
]
add_bullet_list(slide, Inches(0.5), Inches(5.5), Inches(12), Inches(1.5),
                stats_items, size=15, color=ACCENT)

# ════════════════════════════════════════════
# 슬라이드 15: 채택 현실
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
         "에이전틱 AI 채택 현실", size=32, color=WHITE, bold=True)

adoption = [
    ("탐색 중", "30%", 3.0, LIGHT_GRAY),
    ("파일럿", "38%", 3.8, ACCENT2),
    ("배포 준비", "14%", 1.4, ACCENT),
    ("실제 운영", "11%", 1.1, ACCENT3),
]

for i, (label, pct, width_val, color) in enumerate(adoption):
    y = Inches(1.5 + i * 1.1)
    add_text(slide, Inches(0.5), y + Inches(0.05), Inches(2), Inches(0.5),
             label, size=18, color=WHITE)

    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(2.8), y, Inches(width_val), Inches(0.7))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    add_text(slide, Inches(2.8) + Inches(width_val) + Inches(0.2), y + Inches(0.05),
             Inches(1.5), Inches(0.5), pct, size=20, color=color, bold=True)

add_text(slide, Inches(0.5), Inches(5.5), Inches(12), Inches(1),
         "Gartner: 2026년 말까지 엔터프라이즈 앱 40%에 AI 에이전트 내장\n"
         "하지만 40%+ 프로젝트가 레거시 호환 문제로 실패 예측\n"
         "핵심 병목: AI 자체가 아니라 데이터 엔지니어링, 거버넌스, 레거시 통합",
         size=16, color=LIGHT_GRAY)

# ════════════════════════════════════════════
# 슬라이드 16: n8n 소개
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
         "n8n: 에이전틱 AI를 실현하는 플랫폼", size=32, color=N8N_COLOR, bold=True)

# n8n 소개 카드
info_items = [
    ("정체", "Fair-code 워크플로우 자동화 + AI 플랫폼"),
    ("철학", '"코드의 유연성 + 노코드의 속도"'),
    ("통합", "500+ 앱/서비스 연동"),
    ("AI", "LangChain 기반 70+ AI 전용 노드"),
    ("배포", "Self-hosted (무료) 또는 Cloud"),
    ("규모", "$2.5B 밸류에이션, Nvidia 투자"),
]

for i, (key, val) in enumerate(info_items):
    row = i // 2
    col = i % 2
    x = Inches(0.5 + col * 6.3)
    y = Inches(1.3 + row * 1.0)
    add_shape(slide, x, y, Inches(6.0), Inches(0.85), BG_MID)
    add_text(slide, x + Inches(0.3), y + Inches(0.12), Inches(1.5), Inches(0.5),
             key, size=16, color=N8N_COLOR, bold=True)
    add_text(slide, x + Inches(1.8), y + Inches(0.12), Inches(4), Inches(0.5),
             val, size=15, color=WHITE)

# 핵심 숫자들
numbers = [
    ("$180M", "Series C (Nvidia 참여)"),
    ("$40M+", "ARR (전년 대비 10x)"),
    ("3,000+", "엔터프라이즈 고객"),
    ("75%", "고객 AI 도구 활용"),
]
for i, (num, label) in enumerate(numbers):
    x = Inches(0.5 + i * 3.2)
    y = Inches(5.0)
    add_shape(slide, x, y, Inches(2.8), Inches(1.8), BG_MID)
    add_text(slide, x, y + Inches(0.3), Inches(2.8), Inches(0.6),
             num, size=28, color=N8N_COLOR, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x, y + Inches(1.0), Inches(2.8), Inches(0.5),
             label, size=13, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# 슬라이드 17: n8n이 에이전틱 AI 조건을 충족하는 방법
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
         "n8n × 에이전틱 AI 조건 매칭", size=32, color=WHITE, bold=True)

matchings = [
    ("자율성", "트리거 기반 자동 실행\n(Webhook, 스케줄, 이벤트)", ACCENT),
    ("다단계 추론", "AI Agent 노드\nPlan & Execute 패턴", ACCENT2),
    ("도구 사용", "500+ 통합 + MCP\nClient/Server", ACCENT3),
    ("메모리", "Simple / PostgreSQL\n/ Redis 메모리 노드", YELLOW),
    ("적응력", "조건 분기, 에러 핸들링\n재시도 로직", ORANGE),
    ("목표 지향", "워크플로우 = 목표 달성까지의\n전체 파이프라인", N8N_COLOR),
]

for i, (trait, how, color) in enumerate(matchings):
    row = i // 3
    col = i % 3
    x = Inches(0.5 + col * 4.2)
    y = Inches(1.3 + row * 2.8)

    add_shape(slide, x, y, Inches(3.8), Inches(2.4), BG_MID)

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(3.8), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    add_text(slide, x + Inches(0.2), y + Inches(0.3), Inches(3.4), Inches(0.5),
             trait, size=20, color=color, bold=True)
    add_text(slide, x + Inches(0.2), y + Inches(1.0), Inches(3.4), Inches(1.2),
             how, size=15, color=WHITE)

# ════════════════════════════════════════════
# 슬라이드 18: n8n 포지션 맵
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
         "n8n의 독보적 포지션", size=32, color=WHITE, bold=True)

add_text(slide, Inches(0.5), Inches(1.1), Inches(12), Inches(0.5),
         "에이전틱 AI 도구 스펙트럼", size=18, color=LIGHT_GRAY)

# 스펙트럼 바
spectrum = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(0.5), Inches(2.0), Inches(12), Inches(0.6))
spectrum.fill.solid()
spectrum.fill.fore_color.rgb = BG_MID
spectrum.line.fill.background()

add_text(slide, Inches(0.5), Inches(2.0), Inches(3), Inches(0.5),
         "← 코드 전문가용", size=13, color=LIGHT_GRAY)
add_text(slide, Inches(9.5), Inches(2.0), Inches(3), Inches(0.5),
         "비즈니스 팀용 →", size=13, color=LIGHT_GRAY, align=PP_ALIGN.RIGHT)

tools_spectrum = [
    ("LangGraph", Inches(1.2), ACCENT2, "🔴 높음"),
    ("CrewAI", Inches(3.5), ACCENT3, "🟡 중간"),
    ("Agent SDK", Inches(5.5), ACCENT, "🟡 중간"),
    ("n8n", Inches(8.0), N8N_COLOR, "🟢 낮음 ⭐"),
    ("Zapier", Inches(10.5), LIGHT_GRAY, "🟢 낮음"),
]

for name, x, color, difficulty in tools_spectrum:
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, Inches(2.1), Inches(0.4), Inches(0.4))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    add_text(slide, x - Inches(0.5), Inches(2.7), Inches(1.4), Inches(0.4),
             name, size=14, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x - Inches(0.5), Inches(3.1), Inches(1.4), Inches(0.3),
             difficulty, size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# n8n sweet spot 강조
add_text(slide, Inches(2), Inches(3.8), Inches(9), Inches(0.5),
         'n8n = "코드의 유연성 + 노코드의 속도" = Sweet Spot ⭐',
         size=20, color=N8N_COLOR, bold=True, align=PP_ALIGN.CENTER)

# 비교 테이블
comparisons = [
    ("vs LangGraph", "시각적 빌더, 빠른 구축", "복잡한 상태 머신 열세"),
    ("vs CrewAI", "Self-hosted, 데이터 주권", "정교한 크루 구조 열세"),
    ("vs Zapier", "JS/Python 코드 가능, 비용 효율", "통합 수 (500 vs 7,000)"),
]

for i, (vs, pro, con) in enumerate(comparisons):
    y = Inches(4.6 + i * 0.75)
    add_text(slide, Inches(0.5), y, Inches(2.3), Inches(0.5),
             vs, size=15, color=ACCENT, bold=True)
    add_text(slide, Inches(3), y, Inches(4.5), Inches(0.5),
             f"✅ {pro}", size=14, color=ACCENT3)
    add_text(slide, Inches(7.8), y, Inches(4.5), Inches(0.5),
             f"⚠️ {con}", size=14, color=ORANGE)

# ════════════════════════════════════════════
# 슬라이드 19: n8n + MCP
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
         "n8n = 에이전틱 MCP 허브", size=32, color=WHITE, bold=True)

add_text(slide, Inches(0.5), Inches(1.1), Inches(12), Inches(0.5),
         "n8n은 MCP를 양방향으로 지원 (2025.04~)", size=18, color=LIGHT_GRAY)

# 다이어그램
# 왼쪽: MCP 소비
add_shape(slide, Inches(0.5), Inches(2.0), Inches(3.5), Inches(3), BG_MID)
add_text(slide, Inches(0.5), Inches(2.1), Inches(3.5), Inches(0.5),
         "외부 MCP 서버", size=16, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
ext_servers = ["GitHub", "Slack", "PostgreSQL", "파일시스템"]
add_bullet_list(slide, Inches(0.8), Inches(2.7), Inches(3), Inches(2),
                ext_servers, size=14, color=LIGHT_GRAY)

# 가운데: n8n
add_shape(slide, Inches(4.5), Inches(2.0), Inches(4.5), Inches(3), ACCENT2)
add_text(slide, Inches(4.5), Inches(2.2), Inches(4.5), Inches(0.5),
         "n8n 에이전트", size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
n8n_features = [
    "MCP Client → 도구 소비",
    "AI Agent → 오케스트레이션",
    "MCP Server → 도구 노출",
]
add_bullet_list(slide, Inches(4.8), Inches(2.9), Inches(4), Inches(2),
                n8n_features, size=14, color=WHITE)

# 오른쪽: MCP 제공
add_shape(slide, Inches(9.5), Inches(2.0), Inches(3.5), Inches(3), BG_MID)
add_text(slide, Inches(9.5), Inches(2.1), Inches(3.5), Inches(0.5),
         "외부 MCP 클라이언트", size=16, color=N8N_COLOR, bold=True, align=PP_ALIGN.CENTER)
ext_clients = ["Claude Desktop", "Cursor", "VS Code", "다른 에이전트"]
add_bullet_list(slide, Inches(9.8), Inches(2.7), Inches(3), Inches(2),
                ext_clients, size=14, color=LIGHT_GRAY)

# 화살표
add_text(slide, Inches(3.8), Inches(3.0), Inches(0.8), Inches(0.5),
         "→", size=30, color=ACCENT, align=PP_ALIGN.CENTER)
add_text(slide, Inches(8.8), Inches(3.0), Inches(0.8), Inches(0.5),
         "→", size=30, color=N8N_COLOR, align=PP_ALIGN.CENTER)

# 워크플로우 예시
add_text(slide, Inches(0.5), Inches(5.3), Inches(12), Inches(0.5),
         "실제 워크플로우 예시", size=20, color=WHITE, bold=True)

examples = [
    ("고객 문의 에이전트", "이메일 수신 → AI 분류 → 자동 답변/에스컬레이션", ACCENT),
    ("리서치 에이전트", "스케줄 → 웹 검색 → RAG 비교 → 리포트 → Slack 공유", ACCENT3),
    ("DevOps 인시던트", "알림 → 로그 분석 → 원인 추론 → 자동 대응/호출", ORANGE),
]

for i, (name, flow, color) in enumerate(examples):
    y = Inches(5.9 + i * 0.5)
    add_text(slide, Inches(0.8), y, Inches(2.5), Inches(0.4),
             name, size=14, color=color, bold=True)
    add_text(slide, Inches(3.5), y, Inches(9), Inches(0.4),
             flow, size=13, color=LIGHT_GRAY)

# ════════════════════════════════════════════
# 슬라이드 20: n8n 실제 사례
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
         "n8n 실제 도입 사례", size=32, color=WHITE, bold=True)

cases = [
    ("Delivery Hero", "운영 자동화", "월 200시간 절약", ACCENT),
    ("BeGlobal", "AI 맞춤 제안서 생성", "10배 확장, 1분 이내", ACCENT3),
    ("System", "AI 데이터 입력 자동화", "처리 시간 97% 감소", ORANGE),
    ("XIBIX", "HR 질문 AI 워크플로우", "HR 문의 50% 감소", ACCENT2),
    ("Flow AI", "부동산 음성 AI 아웃리치", "수동 후속 완전 자동화", N8N_COLOR),
]

for i, (company, usecase, result, color) in enumerate(cases):
    y = Inches(1.3 + i * 1.15)
    add_shape(slide, Inches(0.5), y, Inches(12), Inches(1.0), BG_MID)

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), y, Inches(0.12), Inches(1.0))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    add_text(slide, Inches(0.9), y + Inches(0.15), Inches(2.5), Inches(0.5),
             company, size=20, color=color, bold=True)
    add_text(slide, Inches(3.5), y + Inches(0.15), Inches(4), Inches(0.5),
             usecase, size=16, color=WHITE)
    add_text(slide, Inches(8), y + Inches(0.15), Inches(4.5), Inches(0.5),
             result, size=18, color=ACCENT3, bold=True)

# ════════════════════════════════════════════
# 슬라이드 21: 솔직한 평가
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
         "솔직한 평가: 가능성과 과제", size=32, color=WHITE, bold=True)

# 가능성
add_shape(slide, Inches(0.5), Inches(1.3), Inches(5.8), Inches(4.5), BG_MID)
add_text(slide, Inches(0.8), Inches(1.4), Inches(5.2), Inches(0.5),
         "✅ 가능성", size=22, color=ACCENT3, bold=True)
pros = [
    "노코드/로우코드로 에이전트 민주화",
    "Self-hosted로 데이터 주권 확보",
    "500+ 통합으로 레거시 시스템 연결",
    "MCP 양방향 지원 → 생태계 허브",
    "시각적 빌더로 빠른 프로토타이핑",
    "Nvidia 투자, $2.5B 밸류에이션",
]
add_bullet_list(slide, Inches(0.8), Inches(2.1), Inches(5.2), Inches(3.5),
                pros, size=15, color=WHITE, spacing=Pt(12))

# 과제
add_shape(slide, Inches(7), Inches(1.3), Inches(5.8), Inches(4.5), BG_MID)
add_text(slide, Inches(7.3), Inches(1.4), Inches(5.2), Inches(0.5),
         "⚠️ 과제", size=22, color=ORANGE, bold=True)
cons = [
    "복잡한 에이전트 로직 → 코드가 유리",
    "엔터프라이즈 거버넌스·보안 검증 필요",
    "LLM 토큰 비용 관리 필요",
    "MCP 자체의 효율성 논쟁",
    "프로덕션 안정성 추가 검증 필요",
    "성공 보장은 없음 (11%만 실운영)",
]
add_bullet_list(slide, Inches(7.3), Inches(2.1), Inches(5.2), Inches(3.5),
                cons, size=15, color=LIGHT_GRAY, spacing=Pt(12))

add_text(slide, Inches(0.5), Inches(6.2), Inches(12), Inches(0.5),
         "성공할지 실패할지는 차치하고 — 시도 자체가 학습이다",
         size=18, color=N8N_COLOR, bold=True, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# 슬라이드 22: 전체 흐름 요약
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
         "전체 흐름 요약", size=32, color=WHITE, bold=True)

flow_steps = [
    ("모델이 똑똑해졌다", "LLM 경쟁", LIGHT_GRAY),
    ("도구를 쓸 수 있게 됐다", "MCP, 네이티브 AI", ACCENT),
    ("코드도 직접 쓴다", "CLI, 바이브 코딩", ACCENT2),
    ("컨텍스트를 관리한다", "컨텍스트 엔지니어링", ACCENT3),
    ("혼자서 일을 한다", "에이전틱 AI", ORANGE),
    ("팀으로 협업한다", "멀티 에이전트", N8N_COLOR),
]

for i, (main_text, sub, color) in enumerate(flow_steps):
    y = Inches(1.2 + i * 0.95)

    add_shape(slide, Inches(2), y, Inches(9), Inches(0.8), BG_MID)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(2), y, Inches(0.1), Inches(0.8))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    add_text(slide, Inches(2.4), y + Inches(0.1), Inches(5), Inches(0.5),
             main_text, size=18, color=WHITE, bold=True)
    add_text(slide, Inches(7.5), y + Inches(0.1), Inches(3), Inches(0.5),
             sub, size=15, color=color)

    if i < 5:
        add_text(slide, Inches(6.3), y + Inches(0.6), Inches(0.5), Inches(0.4),
                 "↓", size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# 결론
add_text(slide, Inches(0.5), Inches(7.0), Inches(12), Inches(0.4),
         "★ 이 모든 것을 시각적으로 조합하고 실행하는 플랫폼 = n8n",
         size=18, color=N8N_COLOR, bold=True, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# 슬라이드 23: 다음 스텝
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.7),
         "다음 스텝 제안", size=32, color=WHITE, bold=True)

next_steps = [
    ("1주차", "n8n Self-hosted 설치\n기본 워크플로우 체험", "도구 익히기", ACCENT),
    ("2주차", "AI Agent 노드로\n간단한 자동화 구축", "가능성 검증", ACCENT2),
    ("3주차", "실제 업무 하나를\n에이전틱 워크플로우로 전환", "실무 적용", ACCENT3),
    ("4주차", "팀 리뷰\n확대 적용 여부 결정", "Go / No-Go", N8N_COLOR),
]

for i, (week, action, goal, color) in enumerate(next_steps):
    x = Inches(0.5 + i * 3.2)
    y = Inches(1.5)

    add_shape(slide, x, y, Inches(2.8), Inches(4.0), BG_MID)

    # 상단 색상 바
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(2.8), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()

    add_text(slide, x, y + Inches(0.3), Inches(2.8), Inches(0.5),
             week, size=24, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.2), y + Inches(1.0), Inches(2.4), Inches(1.5),
             action, size=15, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, x, y + Inches(3.0), Inches(2.8), Inches(0.5),
             goal, size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# 슬라이드 24: 마무리
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
         '"에이전틱 AI의 성패는\n모델이 아니라 오케스트레이션에 달려 있다"',
         size=32, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

# 구분선
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(3.5), Inches(5), Pt(2))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT
line.line.fill.background()

add_text(slide, Inches(1), Inches(4.0), Inches(11), Inches(0.8),
         "성공할지 실패할지는 아직 모른다.\n하지만 해보지 않으면 영원히 모른다.",
         size=22, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.5),
         "에이전틱 AI 시대는 이미 시작되었고,\nn8n은 그 시작을 가장 낮은 문턱에서 열어준다.",
         size=20, color=N8N_COLOR, bold=True, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════
# 슬라이드 25: Q&A
# ════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide, BG_DARK)

add_text(slide, Inches(1), Inches(2.5), Inches(11), Inches(1),
         "Q & A", size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(4.0), Inches(11), Inches(0.5),
         "감사합니다", size=24, color=ACCENT, align=PP_ALIGN.CENTER)

# ── 저장 ──
output = "/Users/mung/git/ppt-maker/AI_트렌드_2025-2026.pptx"
prs.save(output)
print(f"✅ 슬라이드 생성 완료: {output}")
print(f"   총 {len(prs.slides)} 슬라이드")
