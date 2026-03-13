"""AssetManager ↔ SlideBuilder 통합 테스트."""

from __future__ import annotations

import struct
import zlib
from pathlib import Path
from unittest.mock import MagicMock

import pytest
from pptx import Presentation

from ppt_maker.assets.manager import AssetManager
from ppt_maker.converter.hybrid import HybridConverter
from ppt_maker.converter.pptx_custom import SlideBuilder
from ppt_maker.theme.manager import ThemeManager

# --- 헬퍼 ---

def create_test_png(path: Path) -> Path:
    """최소한의 유효한 1x1 PNG 파일 생성."""
    signature = b"\x89PNG\r\n\x1a\n"
    # IHDR: 1x1, 8비트 RGB
    ihdr_data = struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0)
    ihdr_crc = zlib.crc32(b"IHDR" + ihdr_data) & 0xFFFFFFFF
    ihdr = struct.pack(">I", 13) + b"IHDR" + ihdr_data + struct.pack(">I", ihdr_crc)
    # IDAT
    raw_data = zlib.compress(b"\x00\xff\xff\xff")
    idat_crc = zlib.crc32(b"IDAT" + raw_data) & 0xFFFFFFFF
    idat = struct.pack(">I", len(raw_data)) + b"IDAT" + raw_data + struct.pack(">I", idat_crc)
    # IEND
    iend_crc = zlib.crc32(b"IEND") & 0xFFFFFFFF
    iend = struct.pack(">I", 0) + b"IEND" + struct.pack(">I", iend_crc)
    path.write_bytes(signature + ihdr + idat + iend)
    return path


@pytest.fixture()
def theme():
    tm = ThemeManager()
    return tm.load_theme("dark")


@pytest.fixture()
def builder(theme) -> SlideBuilder:
    return SlideBuilder(theme)


@pytest.fixture()
def prs(builder: SlideBuilder) -> Presentation:
    return builder.create_presentation()


# --- SlideBuilder + asset_manager 파라미터 테스트 ---

class TestSlideBuilderAssetManager:
    def test_accepts_asset_manager_parameter(self, theme) -> None:
        """SlideBuilder가 asset_manager 파라미터를 받아야 한다."""
        am = AssetManager()
        sb = SlideBuilder(theme, asset_manager=am)
        assert sb.asset_manager is am

    def test_default_asset_manager_is_none(self, theme) -> None:
        """asset_manager 미지정 시 None이어야 한다."""
        sb = SlideBuilder(theme)
        assert sb.asset_manager is None

    def test_build_picture_left_fallback_no_asset_manager(
        self, builder: SlideBuilder, prs: Presentation
    ) -> None:
        """asset_manager 없으면 플레이스홀더 사각형으로 폴백해야 한다."""
        builder.build_picture_left(prs, {
            "title": "테스트",
            "brand": "BMW",
        })
        assert len(prs.slides) == 1

    def test_build_picture_left_with_real_image(
        self, theme, tmp_path: Path
    ) -> None:
        """asset_manager가 실제 이미지를 반환하면 add_picture로 삽입해야 한다."""
        # 테스트용 PNG 생성
        img_path = create_test_png(tmp_path / "bmw.png")

        # AssetManager를 모킹하여 이미지 경로 반환
        mock_am = MagicMock(spec=AssetManager)
        mock_am.find_brand_image.return_value = img_path

        sb = SlideBuilder(theme, asset_manager=mock_am)
        prs = sb.create_presentation()
        sb.build_picture_left(prs, {"title": "BMW 소개", "brand": "BMW"})

        assert len(prs.slides) == 1
        mock_am.find_brand_image.assert_called_once_with("BMW")
        # 슬라이드에 Picture 도형이 포함되어야 함
        slide = prs.slides[0]
        shape_types = [s.shape_type for s in slide.shapes]
        # shape_type 13 = PICTURE
        assert 13 in shape_types

    def test_build_picture_left_fallback_image_not_found(
        self, theme
    ) -> None:
        """asset_manager가 None을 반환하면 플레이스홀더로 폴백해야 한다."""
        mock_am = MagicMock(spec=AssetManager)
        mock_am.find_brand_image.return_value = None

        sb = SlideBuilder(theme, asset_manager=mock_am)
        prs = sb.create_presentation()
        sb.build_picture_left(prs, {"title": "테스트", "brand": "unknown_brand"})

        assert len(prs.slides) == 1
        mock_am.find_brand_image.assert_called_once_with("unknown_brand")

    def test_build_picture_right_with_real_image(
        self, theme, tmp_path: Path
    ) -> None:
        """build_picture_right도 실제 이미지를 삽입해야 한다."""
        img_path = create_test_png(tmp_path / "mercedes.png")

        mock_am = MagicMock(spec=AssetManager)
        mock_am.find_brand_image.return_value = img_path

        sb = SlideBuilder(theme, asset_manager=mock_am)
        prs = sb.create_presentation()
        sb.build_picture_right(prs, {"title": "메르세데스", "brand": "Mercedes-Benz"})

        assert len(prs.slides) == 1
        mock_am.find_brand_image.assert_called_once_with("Mercedes-Benz")
        slide = prs.slides[0]
        shape_types = [s.shape_type for s in slide.shapes]
        assert 13 in shape_types

    def test_build_picture_right_fallback_no_asset_manager(
        self, builder: SlideBuilder, prs: Presentation
    ) -> None:
        """asset_manager 없으면 우측도 플레이스홀더로 폴백해야 한다."""
        builder.build_picture_right(prs, {
            "title": "성능 비교",
            "content": "최고 출력 544ps",
        })
        assert len(prs.slides) == 1

    def test_build_picture_left_uses_image_key_fallback(
        self, theme, tmp_path: Path
    ) -> None:
        """brand 키가 없으면 image 키를 사용해야 한다."""
        img_path = create_test_png(tmp_path / "audi.png")

        mock_am = MagicMock(spec=AssetManager)
        mock_am.find_brand_image.return_value = img_path

        sb = SlideBuilder(theme, asset_manager=mock_am)
        prs = sb.create_presentation()
        sb.build_picture_left(prs, {"title": "아우디", "image": "Audi"})

        mock_am.find_brand_image.assert_called_once_with("Audi")


# --- HybridConverter asset_manager 전달 테스트 ---

class TestHybridConverterAssetManager:
    def test_passes_asset_manager_to_builder(self, theme) -> None:
        """HybridConverter가 asset_manager를 SlideBuilder에 전달해야 한다."""
        am = AssetManager()
        conv = HybridConverter(theme, asset_manager=am)
        assert conv.builder.asset_manager is am

    def test_default_asset_manager_is_none(self, theme) -> None:
        """asset_manager 미지정 시 builder도 None이어야 한다."""
        conv = HybridConverter(theme)
        assert conv.builder.asset_manager is None
