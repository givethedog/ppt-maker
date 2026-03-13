"""template/analyzer 모듈 테스트."""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation

from ppt_maker.errors import ConfigError
from ppt_maker.template.analyzer import (
    TemplateManifest,
    _auto_map_layout,
    _emu_to_inches,
    analyze_template,
    load_manifest,
    save_manifest,
)


def _create_test_pptx(path: Path) -> Path:
    """테스트용 .pptx 파일 생성."""
    prs = Presentation()
    # 빈 프레젠테이션은 기본 레이아웃을 가짐
    layout = prs.slide_layouts[0]  # Title Slide
    prs.slides.add_slide(layout)
    prs.save(str(path))
    return path


def test_emu_to_inches():
    assert _emu_to_inches(914400) == 1.0
    assert _emu_to_inches(0) == 0.0


def test_auto_map_layout():
    assert _auto_map_layout("Title Slide") == "title"
    assert _auto_map_layout("제목 슬라이드") == "title"
    assert _auto_map_layout("Section Header") == "section"
    assert _auto_map_layout("Title and Content") == "content"
    assert _auto_map_layout("Two Content") == "two_content"
    assert _auto_map_layout("Blank") == "blank"
    assert _auto_map_layout("Custom Layout XYZ") is None
    # BMW CI 레이아웃 이름
    assert _auto_map_layout("Title | Full Area") == "title"
    assert _auto_map_layout("Divider | Half Area Right") == "section"
    assert _auto_map_layout("Content | 1") == "content"
    assert _auto_map_layout("Content | 2") == "two_content"
    assert _auto_map_layout("Content | 2x2") == "comparison"
    assert _auto_map_layout("Grid | 2x2") == "grid_2x2"
    assert _auto_map_layout("Content | Area") == "content_area"
    assert _auto_map_layout("Content | Area DARK") == "content_area_dark"
    assert _auto_map_layout("Content | Picture Left") == "picture_left"
    assert _auto_map_layout("Content | Picture Right") == "picture_right"
    assert _auto_map_layout("Key Note") == "keynote"


def test_analyze_template_nonexistent():
    with pytest.raises(ConfigError, match="찾을 수 없습니다"):
        analyze_template(Path("/nonexistent/template.pptx"))


def test_analyze_template(tmp_path):
    pptx_path = tmp_path / "test.pptx"
    _create_test_pptx(pptx_path)

    manifest = analyze_template(pptx_path)

    assert manifest.source_path == str(pptx_path)
    assert manifest.slide_width > 0
    assert manifest.slide_height > 0
    assert len(manifest.layouts) > 0


def test_analyze_template_layouts(tmp_path):
    pptx_path = tmp_path / "test.pptx"
    _create_test_pptx(pptx_path)

    manifest = analyze_template(pptx_path)

    # python-pptx 기본 프레젠테이션은 여러 레이아웃을 가짐
    layout_names = [lay.name for lay in manifest.layouts]
    assert len(layout_names) > 0
    # 각 레이아웃에 인덱스가 있어야 함
    for i, layout in enumerate(manifest.layouts):
        assert layout.index == i


def test_save_and_load_manifest(tmp_path):
    manifest = TemplateManifest(
        source_path="/test/template.pptx",
        slide_width=13.333,
        slide_height=7.5,
        fonts_used={"Arial", "맑은 고딕"},
        layout_mapping={"title": 0, "content": 1},
    )

    output_path = tmp_path / "manifest.toml"
    save_manifest(manifest, output_path)

    assert output_path.exists()

    loaded = load_manifest(output_path)
    assert loaded.source_path == "/test/template.pptx"
    assert loaded.slide_width == 13.333
    assert loaded.slide_height == 7.5
    assert loaded.fonts_used == {"Arial", "맑은 고딕"}
    assert loaded.layout_mapping == {"title": 0, "content": 1}


def test_load_manifest_nonexistent():
    with pytest.raises(ConfigError, match="매니페스트"):
        load_manifest(Path("/nonexistent/manifest.toml"))


def test_manifest_summary():
    from ppt_maker.template.analyzer import LayoutInfo

    manifest = TemplateManifest(
        source_path="test.pptx",
        slide_width=10.0,
        slide_height=7.5,
        layouts=[LayoutInfo(name="Title Slide", index=0)],
        layout_mapping={"title": 0},
    )
    summary = manifest.summary()
    assert "test.pptx" in summary
    assert "10.0 x 7.5" in summary
    assert "Title Slide" in summary


def test_analyze_real_template(tmp_path):
    """BMW CI Template이 있으면 실제 분석 테스트."""
    real_template = Path("examples/simple/BMW CI Template.pptx")
    if not real_template.exists():
        pytest.skip("BMW CI Template not found")

    manifest = analyze_template(real_template)
    assert len(manifest.layouts) > 0
    assert manifest.slide_width > 0

    # 매니페스트 저장/로드 라운드트립
    output = tmp_path / "bmw_manifest.toml"
    save_manifest(manifest, output)
    loaded = load_manifest(output)
    assert len(loaded.layouts) == len(manifest.layouts)
