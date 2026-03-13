"""pptx_custom.py 단위 테스트."""

from __future__ import annotations

from pathlib import Path

import pytest
from pptx import Presentation

from ppt_maker.converter.pptx_custom import SlideBuilder
from ppt_maker.errors import PptxCustomError
from ppt_maker.theme.manager import ThemeManager


@pytest.fixture()
def builder() -> SlideBuilder:
    tm = ThemeManager()
    theme = tm.load_theme("dark")
    return SlideBuilder(theme)


@pytest.fixture()
def prs(builder: SlideBuilder) -> Presentation:
    return builder.create_presentation()


class TestSlideBuilder:
    def test_build_title(self, builder: SlideBuilder, prs: Presentation) -> None:
        builder.build_title(prs, {"title": "테스트 제목", "subtitle": "부제", "date": "2026-03"})
        assert len(prs.slides) == 1

    def test_build_quote(self, builder: SlideBuilder, prs: Presentation) -> None:
        builder.build_quote(prs, {"quote": "AI는 미래다", "attribution": "전문가"})
        assert len(prs.slides) == 1

    def test_build_timeline(self, builder: SlideBuilder, prs: Presentation) -> None:
        builder.build_timeline(prs, {
            "title": "AI 발전 타임라인",
            "events": [
                {"date": "2023", "title": "GPT-4 출시"},
                {"date": "2024", "title": "Claude 3 출시"},
                {"date": "2025", "title": "에이전틱 AI"},
            ],
        })
        assert len(prs.slides) == 1

    def test_build_comparison(self, builder: SlideBuilder, prs: Presentation) -> None:
        builder.build_comparison(prs, {
            "title": "비교",
            "left": {"title": "전통 방식", "items": ["수동 작업", "느림"]},
            "right": {"title": "AI 방식", "items": ["자동화", "빠름"]},
        })
        assert len(prs.slides) == 1

    def test_build_card_list(self, builder: SlideBuilder, prs: Presentation) -> None:
        builder.build_card_list(prs, {
            "title": "주요 기능",
            "items": [
                {"title": "기능 1", "detail": "설명 1"},
                {"title": "기능 2", "detail": "설명 2"},
            ],
        })
        assert len(prs.slides) == 1

    def test_build_process(self, builder: SlideBuilder, prs: Presentation) -> None:
        builder.build_process(prs, {
            "title": "파이프라인",
            "steps": [
                {"label": "입력", "desc": "주제 입력"},
                {"label": "변환", "desc": "MD→PPTX"},
                {"label": "출력", "desc": "파일 생성"},
            ],
        })
        assert len(prs.slides) == 1

    def test_unsupported_type_raises(self, builder: SlideBuilder, prs: Presentation) -> None:
        with pytest.raises(PptxCustomError, match="지원되지 않는"):
            builder.build_slide(prs, "nonexistent_type", {})

    def test_build_slide_dispatcher(self, builder: SlideBuilder, prs: Presentation) -> None:
        builder.build_slide(prs, "title", {"title": "디스패처 테스트"})
        assert len(prs.slides) == 1

    def test_build_section(self, builder: SlideBuilder, prs: Presentation) -> None:
        builder.build_section(prs, {"title": "제2장", "subtitle": "AI의 미래"})
        assert len(prs.slides) == 1

    def test_build_picture_left(self, builder: SlideBuilder, prs: Presentation) -> None:
        builder.build_picture_left(prs, {
            "title": "차량 소개",
            "items": ["특징 1", "특징 2"],
            "image_label": "[BMW i4]",
        })
        assert len(prs.slides) == 1

    def test_build_picture_right(self, builder: SlideBuilder, prs: Presentation) -> None:
        builder.build_picture_right(prs, {
            "title": "성능 비교",
            "content": "최고 출력 544ps",
        })
        assert len(prs.slides) == 1

    def test_build_blank(self, builder: SlideBuilder, prs: Presentation) -> None:
        builder.build_blank(prs, {"title": "빈 슬라이드"})
        assert len(prs.slides) == 1

    def test_build_blank_dark(self, builder: SlideBuilder, prs: Presentation) -> None:
        builder.build_blank_dark(prs, {"title": "다크 슬라이드"})
        assert len(prs.slides) == 1

    def test_build_keynote(self, builder: SlideBuilder, prs: Presentation) -> None:
        builder.build_keynote(prs, {
            "title": "핵심 메시지",
            "message": "혁신은 계속됩니다",
        })
        assert len(prs.slides) == 1

    def test_layout_mapping(self) -> None:
        """레이아웃 매핑이 적용되는지 확인."""
        tm = ThemeManager()
        theme = tm.load_theme("dark")
        mapping = {"title": 0, "content": 1}
        builder = SlideBuilder(theme, layout_mapping=mapping)
        assert builder._layout_mapping == mapping

    def test_save_pptx(self, builder: SlideBuilder, prs: Presentation, tmp_path: Path) -> None:
        builder.build_title(prs, {"title": "저장 테스트"})
        builder.build_timeline(prs, {"title": "TL", "events": [{"date": "2025", "title": "E"}]})
        out = tmp_path / "test.pptx"
        prs.save(str(out))
        assert out.exists()
        loaded = Presentation(str(out))
        assert len(loaded.slides) == 2

    def test_create_presentation_with_template(self, tmp_path: Path) -> None:
        """회사 템플릿 기반 프레젠테이션 생성."""
        real_template = Path("examples/simple/BMW CI Template.pptx")
        if not real_template.exists():
            pytest.skip("BMW CI Template not found")

        tm = ThemeManager()
        theme = tm.load_theme("dark")
        mapping = {"title": 0, "section": 1, "content": 7}
        builder = SlideBuilder(theme, layout_mapping=mapping)
        prs = builder.create_presentation(template_path=real_template)

        # 기존 슬라이드는 제거되어야 함
        assert len(prs.slides) == 0
        # 레이아웃은 유지되어야 함
        assert len(prs.slide_layouts) == 20

        # 슬라이드 추가 후 저장
        builder.build_title(prs, {"title": "BMW IT Report"})
        builder.build_section(prs, {"title": "Chapter 1"})
        assert len(prs.slides) == 2

        out = tmp_path / "bmw_test.pptx"
        prs.save(str(out))
        assert out.exists()
